"""
PPTX Slide Editor — Flask Backend
Edits NotebookLM image-based slides with text/shape overlays,
then exports a new editable PPTX.
"""

import os, sys, json, base64, subprocess, shutil, uuid, datetime, re
from pathlib import Path
from flask import Flask, render_template, jsonify, request, send_file, make_response
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image, ImageDraw, ImageFont
from werkzeug.utils import secure_filename
import io
import threading

# Cap PIL image dimensions to prevent decompression-bomb DoS
Image.MAX_IMAGE_PIXELS = 50_000_000  # ~50 MP

# Strict name pattern for templates/version snapshots (defence in depth)
_SAFE_NAME_RE = re.compile(r"^[A-Za-z0-9_\- ]{1,64}$")

def _safe_name(s):
    """Return a safe name or empty string if invalid. Use for templates/versions."""
    if not isinstance(s, str):
        return ""
    s = s.strip()
    if not s or not _SAFE_NAME_RE.match(s):
        return ""
    if s in (".", ".."):
        return ""
    return s

def _ensure_dict(payload):
    """Return payload if dict, else empty dict."""
    return payload if isinstance(payload, dict) else {}

# OCR — use EasyOCR (no external binary needed, works out of the box)
HAS_OCR = False
_ocr_reader = None
try:
    import easyocr
    HAS_OCR = True
except ImportError:
    pass

# Background removal — rembg (optional, ~170 MB model download on first use)
HAS_REMBG = False
try:
    from rembg import remove as rembg_remove
    HAS_REMBG = True
except ImportError:
    pass

app = Flask(__name__)
# Default upload cap is 1 GB so the bulk endpoint comfortably accepts 20 PPTX
# files even when individual files run heavy (~30-50 MB with embedded video).
# For a stricter single-file-only deployment, set MAX_UPLOAD_MB=60 in the env.
app.config['MAX_CONTENT_LENGTH'] = int(os.environ.get('MAX_UPLOAD_MB', 1024)) * 1024 * 1024
# Cap base64 image-overlay payloads embedded in JSON requests (bytes after decode)
MAX_OVERLAY_IMG_BYTES = 8 * 1024 * 1024  # 8 MB decoded


# Friendly 413 handler so the bulk uploader can show a useful toast instead of
# Flask's default HTML page.
from werkzeug.exceptions import RequestEntityTooLarge


@app.errorhandler(RequestEntityTooLarge)
def _too_large(e):
    cap_mb = app.config['MAX_CONTENT_LENGTH'] // (1024 * 1024)
    return jsonify({
        "error": (
            f"Upload too large — combined size exceeds {cap_mb} MB cap. "
            f"Re-run with MAX_UPLOAD_MB=500 (or larger) to allow it, "
            f"or split your batch into smaller groups."
        ),
        "cap_mb": cap_mb,
    }), 413

BASE_DIR   = Path(__file__).parent
SLIDES_DIR = BASE_DIR / "static" / "slides"
ORIGINALS_DIR = SLIDES_DIR / "_originals"
UPLOAD_DIR = BASE_DIR / "uploads"
EXPORT_DIR = BASE_DIR / "exports"
DATA_FILE  = BASE_DIR / "slide_data.json"
VIDEO_DIR  = BASE_DIR / "videos"
EXPORT_DIR.mkdir(exist_ok=True)
VIDEO_DIR.mkdir(exist_ok=True)
ORIGINALS_DIR.mkdir(parents=True, exist_ok=True)

# Max age (seconds) for files in EXPORT_DIR. Older files purged on next export.
EXPORT_TTL_SECONDS = int(os.environ.get('EXPORT_TTL_SECONDS', 24 * 3600))

def _cleanup_old_exports():
    """Best-effort cleanup of EXPORT_DIR files older than EXPORT_TTL_SECONDS."""
    import time as _t
    cutoff = _t.time() - EXPORT_TTL_SECONDS
    try:
        for f in EXPORT_DIR.iterdir():
            try:
                if f.is_file() and f.stat().st_mtime < cutoff:
                    f.unlink()
            except OSError:
                pass
    except OSError:
        pass

SLIDE_W_PX, SLIDE_H_PX = 2134, 1200   # actual JPG pixel dimensions

# Lock for file-based persistence (slide_data.json, comments.json)
# RLock (reentrant) so nested critical sections — e.g. ops_undo holds the lock
# while calling _restore_from_snapshot, which also acquires it — don't deadlock.
_data_lock = threading.RLock()

ALLOWED_VIDEO_EXTS = {'.mp4', '.mov', '.avi', '.mkv', '.webm', '.flv', '.wmv'}
ALLOWED_IMAGE_FORMATS = {'PNG', 'JPEG', 'GIF', 'WEBP', 'BMP'}

IS_WINDOWS = sys.platform == "win32"
IS_MACOS = sys.platform == "darwin"
IS_LINUX = sys.platform.startswith("linux")

def _get_slide_files():
    return sorted(SLIDES_DIR.glob("slide-*.jpg"))


def _parse_hex_color(s, default=(0, 0, 0)):
    """Parse a hex color string like '#FF00AA' or 'FF00AA'. Returns (r, g, b) tuple."""
    try:
        s = s.lstrip("#")
        if len(s) < 6:
            return default
        return (int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))
    except (ValueError, TypeError, AttributeError):
        return default


# ── Persistence ──────────────────────────────────────────────────────────────

def load_data():
    with _data_lock:
        if DATA_FILE.exists():
            return json.loads(DATA_FILE.read_text())
        num_slides = len(_get_slide_files())
        return {str(i+1): {"overlays": [], "notes": ""} for i in range(num_slides)}

def save_data(data):
    with _data_lock:
        DATA_FILE.write_text(json.dumps(data, indent=2))

# ── Active deck name ────────────────────────────────────────────────────────

DECK_NAME_FILE = BASE_DIR / "current_deck.txt"

def _set_deck_name(name):
    try:
        DECK_NAME_FILE.write_text((name or "").strip()[:255])
    except OSError:
        pass

def _get_deck_name():
    try:
        if DECK_NAME_FILE.exists():
            return DECK_NAME_FILE.read_text().strip()
    except OSError:
        pass
    return ""

def _parse_skip_slides(value):
    """Normalize a 'skip slides' input into a set of 1-based ints.

    Accepts:
      - list/tuple of ints or numeric strings
      - comma-separated string with optional ranges, e.g. '2,4,7-9'
    Invalid tokens are silently dropped. Returns set() if nothing valid.
    """
    out = set()
    if value is None:
        return out
    tokens = []
    if isinstance(value, (list, tuple)):
        tokens = list(value)
    elif isinstance(value, str):
        tokens = [t.strip() for t in value.split(',') if t.strip()]
    else:
        return out
    for tok in tokens:
        if isinstance(tok, int):
            if tok >= 1:
                out.add(tok)
            continue
        if not isinstance(tok, str):
            continue
        tok = tok.strip()
        if not tok:
            continue
        if '-' in tok:
            parts = tok.split('-', 1)
            try:
                a, b = int(parts[0].strip()), int(parts[1].strip())
            except ValueError:
                continue
            if a > b:
                a, b = b, a
            for i in range(max(1, a), b + 1):
                out.add(i)
        else:
            try:
                n = int(tok)
            except ValueError:
                continue
            if n >= 1:
                out.add(n)
    return out


# ── NotebookLM Logo Removal ─────────────────────────────────────────────────

# Logo region at 1376×768 source resolution: bottom-right 140×25 px.
# We scale proportionally to whatever the actual image size is.
LOGO_REF_W, LOGO_REF_H = 1376, 768   # reference image dimensions
LOGO_W, LOGO_H          = 145, 28    # tight box around icon + "NotebookLM" text (with small pad)

def _erase_logo(img):
    """Erase only the NotebookLM logo by tiling the pixel row just above it
    downward. Preserves the exact horizontal background pattern — pixel-perfect."""
    w, h = img.size
    logo_w = int(LOGO_W * w / LOGO_REF_W)
    logo_h = int(LOGO_H * h / LOGO_REF_H)

    # Single-row strip just above the logo (1 px tall)
    src_y = max(0, h - logo_h - 1)
    strip = img.crop((w - logo_w, src_y, w, src_y + 1))

    # Tile that row to full logo height (NEAREST = exact pixel copy)
    patch = strip.resize((logo_w, logo_h), Image.NEAREST)
    img.paste(patch, (w - logo_w, h - logo_h))
    return img


def remove_logos_batch(slide_files):
    """Remove the logo from many slides with minimal per-file overhead."""
    for p in slide_files:
        img = Image.open(p).convert("RGB")
        img = _erase_logo(img)
        img.save(str(p), quality=95)


def process_uploaded_pptx(pptx_path):
    """Convert PPTX slides to JPG images. Atomic: stages to a temp dir and
    only swaps into place on success. Originals are preserved for reset."""
    import tempfile
    SLIDES_DIR.mkdir(parents=True, exist_ok=True)
    stage_dir = Path(tempfile.mkdtemp(prefix="slides_stage_"))
    try:
        _convert_pptx_to_images_libreoffice(pptx_path, stage_dir)

        staged = sorted(stage_dir.glob("slide-*.jpg"))
        if not staged:
            raise RuntimeError("Conversion produced no slide images")

        # Conversion succeeded — now swap. Wipe live + originals dirs.
        for f in SLIDES_DIR.glob("slide-*.jpg"):
            f.unlink()
        for f in ORIGINALS_DIR.glob("slide-*.jpg"):
            f.unlink()
        if DATA_FILE.exists():
            DATA_FILE.unlink()

        for sf in staged:
            shutil.copy2(str(sf), str(ORIGINALS_DIR / sf.name))
            shutil.move(str(sf), str(SLIDES_DIR / sf.name))
    finally:
        shutil.rmtree(stage_dir, ignore_errors=True)


def _find_libreoffice():
    """Return the LibreOffice binary path or None (cross-platform)."""
    candidates = ["libreoffice", "soffice"]
    if IS_WINDOWS:
        candidates.extend([
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
            str(Path.home() / "AppData" / "Local" / "Programs" / "LibreOffice" / "program" / "soffice.exe"),
        ])
    elif IS_MACOS:
        candidates.append("/Applications/LibreOffice.app/Contents/MacOS/soffice")
    else:  # Linux
        candidates.extend([
            "/usr/bin/soffice",
            "/snap/bin/libreoffice",
            "/usr/bin/libreoffice",
        ])
    for candidate in candidates:
        if shutil.which(candidate) or Path(candidate).exists():
            return candidate
    return None


def _convert_pptx_to_images_libreoffice(pptx_path, output_dir=None):
    """PPTX → PDF (LibreOffice) → JPG per page (pdf2image or PyMuPDF)."""
    import tempfile
    dest = output_dir or SLIDES_DIR
    lo_cmd = _find_libreoffice()
    if not lo_cmd:
        raise RuntimeError("LibreOffice not found")

    tmp_dir = Path(tempfile.mkdtemp())
    try:
        subprocess.run(
            [lo_cmd, "--headless", "--convert-to", "pdf",
             "--outdir", str(tmp_dir), str(pptx_path)],
            check=True, timeout=120,
        )
        pdf_files = list(tmp_dir.glob("*.pdf"))
        if not pdf_files:
            raise RuntimeError("PDF conversion produced no output")

        try:
            import fitz
            doc = fitz.open(str(pdf_files[0]))
            mat = fitz.Matrix(2.0, 2.0)
            for i, page in enumerate(doc):
                pix = page.get_pixmap(matrix=mat)
                (dest / f"slide-{i+1:02d}.jpg").write_bytes(pix.tobytes("jpeg"))
            doc.close()
        except ImportError:
            from pdf2image import convert_from_path
            for i, img in enumerate(convert_from_path(str(pdf_files[0]), dpi=200)):
                img.convert("RGB").save(
                    str(dest / f"slide-{i+1:02d}.jpg"), "JPEG", quality=95)
    finally:
        shutil.rmtree(tmp_dir, ignore_errors=True)


# ── Routes ───────────────────────────────────────────────────────────────────

@app.route("/")
def index():
    slide_files = _get_slide_files()
    slides = [{"index": i+1, "file": f"slides/{f.name}"} for i, f in enumerate(slide_files)]
    resp = make_response(render_template("index.html", slides=slides, num_slides=len(slide_files)))
    resp.headers["Cache-Control"] = "no-cache, no-store, must-revalidate"
    resp.headers["Pragma"] = "no-cache"
    return resp

@app.route("/api/slide/<int:num>")
def get_slide(num):
    data = load_data()
    return jsonify(data.get(str(num), {"overlays": [], "notes": ""}))

@app.route("/api/upload", methods=["POST"])
def upload_pptx():
    """Upload a PPTX file, convert to slide images, and remove NotebookLM logo."""
    if "file" not in request.files:
        return jsonify({"error": "No file uploaded"}), 400
    f = request.files["file"]
    if not f.filename.lower().endswith(".pptx"):
        return jsonify({"error": "Only .pptx files are supported"}), 400

    UPLOAD_DIR.mkdir(exist_ok=True)
    save_path = UPLOAD_DIR / secure_filename(f.filename)
    f.save(str(save_path))

    try:
        process_uploaded_pptx(save_path)
    except Exception as e:
        return jsonify({"error": f"Processing failed: {e}"}), 500

    _set_deck_name(f.filename)
    return jsonify({"ok": True, "num_slides": len(_get_slide_files()),
                    "deck_name": _get_deck_name()})


@app.route("/api/remove-logo", methods=["POST"])
def remove_logo_from_existing():
    """Remove NotebookLM logo from all currently loaded slide images."""
    slides = sorted(SLIDES_DIR.glob("slide-*.jpg"))
    if not slides:
        return jsonify({"ok": True, "count": 0})
    snapshot = _snapshot_before_destructive("remove-logo")
    remove_logos_batch(slides)
    log_id = _log_op("remove-logo",
                     text=f"Removed NotebookLM logo from {len(slides)} slide(s)",
                     scope="all", count=len(slides), snapshot=snapshot)
    return jsonify({"ok": True, "count": len(slides),
                    "snapshot": snapshot, "log_id": log_id})


@app.route("/api/slide/<int:num>/reset", methods=["POST"])
def reset_slide(num):
    """Restore a slide from its original (kept on upload). Undoes destructive edits."""
    slide_files = _get_slide_files()
    if num < 1 or num > len(slide_files):
        return jsonify({"error": "Invalid slide"}), 400
    target = slide_files[num - 1]
    orig = ORIGINALS_DIR / target.name
    if not orig.exists():
        return jsonify({"error": "No original available for this slide"}), 404
    shutil.copy2(str(orig), str(target))
    # Clear overlays for this slide
    with _data_lock:
        data = json.loads(DATA_FILE.read_text()) if DATA_FILE.exists() else {}
        data[str(num)] = {"overlays": [], "notes": data.get(str(num), {}).get("notes", "")}
        DATA_FILE.write_text(json.dumps(data, indent=2))
    return jsonify({"ok": True})


@app.route("/api/slide/<int:num>/inpaint-region", methods=["POST"])
def inpaint_region(num):
    """Inpaint (erase) a normalised {x,y,w,h} region of a slide using
    cv2.INPAINT_TELEA. Used by the OCR edit flow to remove the original text
    pixels while preserving the surrounding background gradient/texture."""
    import cv2
    import numpy as np
    slide_files = _get_slide_files()
    if num < 1 or num > len(slide_files):
        return jsonify({"error": "Invalid slide"}), 400
    payload = _ensure_dict(request.json)
    try:
        rx = float(payload.get("x", 0))
        ry = float(payload.get("y", 0))
        rw = float(payload.get("w", 0))
        rh = float(payload.get("h", 0))
    except (TypeError, ValueError):
        return jsonify({"error": "x/y/w/h must be numeric"}), 400
    if rw <= 0 or rh <= 0:
        return jsonify({"error": "Region width/height must be positive"}), 400

    snapshot = _snapshot_before_destructive("inpaint-region")
    path = slide_files[num - 1]
    img = cv2.imread(str(path))
    if img is None:
        return jsonify({"error": "Could not read slide image"}), 500
    h, w = img.shape[:2]
    # Slight outward pad so anti-aliased edges of the original glyphs are erased too
    pad = max(2, int(min(w, h) * 0.004))
    x1 = max(0, int(rx * w) - pad)
    y1 = max(0, int(ry * h) - pad)
    x2 = min(w, int((rx + rw) * w) + pad)
    y2 = min(h, int((ry + rh) * h) + pad)
    mask = np.zeros((h, w), dtype=np.uint8)
    mask[y1:y2, x1:x2] = 255
    result = cv2.inpaint(img, mask, inpaintRadius=5, flags=cv2.INPAINT_TELEA)
    cv2.imwrite(str(path), result, [cv2.IMWRITE_JPEG_QUALITY, 95])
    log_id = _log_op("inpaint-region",
                     text=f"Erased region on slide {num} for OCR edit",
                     scope="current", slide_num=num, snapshot=snapshot)
    return jsonify({"ok": True, "snapshot": snapshot, "log_id": log_id})


@app.route("/api/slide/<int:num>", methods=["POST"])
def save_slide(num):
    num_slides = len(_get_slide_files())
    if num < 1 or num > max(num_slides, 1):
        return jsonify({"error": "Invalid slide number"}), 400
    payload = request.json
    if not isinstance(payload, dict):
        return jsonify({"error": "Invalid payload"}), 400
    # Validate expected keys
    if "overlays" in payload and not isinstance(payload["overlays"], list):
        return jsonify({"error": "overlays must be a list"}), 400
    if "notes" in payload and not isinstance(payload["notes"], str):
        return jsonify({"error": "notes must be a string"}), 400
    # Whitelist keys to prevent arbitrary data injection
    clean = {
        "overlays": payload.get("overlays", []),
        "notes": payload.get("notes", ""),
    }
    with _data_lock:
        data = json.loads(DATA_FILE.read_text()) if DATA_FILE.exists() else {}
        data[str(num)] = clean
        DATA_FILE.write_text(json.dumps(data, indent=2))
    return jsonify({"ok": True})

@app.route("/api/export", methods=["POST"])
def export_pptx():
    """Rebuild PPTX: slide image as background + text overlays as real text boxes."""
    _cleanup_old_exports()
    data = load_data()
    prs  = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # completely blank

    current_slides = _get_slide_files()
    for i, slide_file in enumerate(current_slides):
        slide_num = str(i + 1)
        slide     = prs.slides.add_slide(blank_layout)

        # Background image (full slide)
        pic = slide.shapes.add_picture(
            str(slide_file),
            left=0, top=0,
            width=prs.slide_width, height=prs.slide_height
        )
        # Send image to back
        slide.shapes._spTree.remove(pic._element)
        slide.shapes._spTree.insert(2, pic._element)

        # Add overlays
        slide_data = data.get(slide_num, {})
        for ov in slide_data.get("overlays", []):
            _add_overlay(slide, ov, prs.slide_width, prs.slide_height)

        # Speaker notes
        notes = slide_data.get("notes", "").strip()
        if notes:
            tf = slide.notes_slide.notes_text_frame
            tf.text = notes

    export_name = f"SlideCraft_Export_{uuid.uuid4().hex[:8]}.pptx"
    out_path = EXPORT_DIR / export_name
    prs.save(str(out_path))
    return send_file(str(out_path), as_attachment=True,
                     download_name="SlideCraft_Export.pptx")

@app.route("/api/slide/<int:num>/preview", methods=["POST"])
def preview_slide(num):
    """Render composite preview: slide image + overlays baked in."""
    slide_files = _get_slide_files()
    if num < 1 or num > len(slide_files):
        return jsonify({"error": "Invalid slide"}), 400

    data   = load_data()
    payload = request.json
    data[str(num)] = payload
    save_data(data)

    slide_file = slide_files[num - 1]
    img = Image.open(slide_file).convert("RGBA")
    draw = ImageDraw.Draw(img)

    for ov in payload.get("overlays", []):
        x = int(ov["x"] * SLIDE_W_PX)
        y = int(ov["y"] * SLIDE_H_PX)
        w = int(ov["w"] * SLIDE_W_PX)
        h = int(ov["h"] * SLIDE_H_PX)

        if ov["type"] == "text":
            # Semi-transparent bg box
            bg = ov.get("bgColor", "transparent")
            if bg and bg != "transparent":
                overlay_img = Image.new("RGBA", img.size, (0,0,0,0))
                ov_draw = ImageDraw.Draw(overlay_img)
                cr, cg, cb = _parse_hex_color(bg)
                ov_draw.rectangle([x, y, x+w, y+h], fill=(cr,cg,cb,180))
                img = Image.alpha_composite(img, overlay_img)
                draw = ImageDraw.Draw(img)

            r, g, b = _parse_hex_color(ov.get("color", "#FFFFFF"), (255, 255, 255))
            font_size = int(ov.get("fontSize", 18) * 2.5)
            font = _get_bake_font(font_size, bold=True, family=ov.get("fontFamily", "Arial"))
            draw.text((x+8, y+8), ov.get("text",""), fill=(r,g,b,255), font=font)

        elif ov["type"] == "rect":
            r, g, b = _parse_hex_color(ov.get("fillColor", "#2563EB"), (37, 99, 235))
            overlay_img = Image.new("RGBA", img.size, (0,0,0,0))
            ov_draw = ImageDraw.Draw(overlay_img)
            ov_draw.rectangle([x,y,x+w,y+h], fill=(r,g,b,160))
            img = Image.alpha_composite(img, overlay_img)
            draw = ImageDraw.Draw(img)

    img_rgb = img.convert("RGB")
    buf = io.BytesIO()
    img_rgb.save(buf, format="JPEG", quality=88)
    buf.seek(0)
    b64 = base64.b64encode(buf.read()).decode()
    return jsonify({"preview": f"data:image/jpeg;base64,{b64}"})


# ── Cross-platform font helpers ────────────────────────────────────────────

def _get_font_dirs():
    """Return a list of directories where system fonts are installed."""
    dirs = []
    if IS_WINDOWS:
        win_fonts = Path(os.environ.get("WINDIR", r"C:\Windows")) / "Fonts"
        dirs.append(win_fonts)
        # User-installed fonts on Windows 10+
        local_fonts = Path.home() / "AppData" / "Local" / "Microsoft" / "Windows" / "Fonts"
        if local_fonts.exists():
            dirs.append(local_fonts)
    elif IS_MACOS:
        dirs.extend([
            Path("/Library/Fonts"),
            Path("/System/Library/Fonts"),
            Path("/System/Library/Fonts/Supplemental"),
            Path.home() / "Library" / "Fonts",
        ])
    else:  # Linux and others
        dirs.extend([
            Path("/usr/share/fonts"),
            Path("/usr/local/share/fonts"),
            Path.home() / ".fonts",
            Path.home() / ".local" / "share" / "fonts",
        ])
    return [d for d in dirs if d.exists()]


# Font family → {(bold, italic): [possible filenames]} mapping
# Filenames only — resolved against font dirs at runtime
_FONT_FILENAMES = {
    "Arial": {
        (False, False): ["arial.ttf", "Arial.ttf"],
        (True, False):  ["arialbd.ttf", "Arial Bold.ttf", "Arial-BoldMT.ttf"],
        (False, True):  ["ariali.ttf", "Arial Italic.ttf"],
        (True, True):   ["arialbi.ttf", "Arial Bold Italic.ttf"],
    },
    "Segoe UI": {
        (False, False): ["segoeui.ttf", "SegoeUI.ttf"],
        (True, False):  ["segoeuib.ttf", "SegoeUI-Bold.ttf"],
        (False, True):  ["segoeuii.ttf", "SegoeUI-Italic.ttf"],
        (True, True):   ["segoeuiz.ttf", "SegoeUI-BoldItalic.ttf"],
    },
    "Calibri": {
        (False, False): ["calibri.ttf", "Calibri.ttf"],
        (True, False):  ["calibrib.ttf", "Calibri-Bold.ttf"],
        (False, True):  ["calibrii.ttf", "Calibri-Italic.ttf"],
        (True, True):   ["calibriz.ttf", "Calibri-BoldItalic.ttf"],
    },
    "Verdana": {
        (False, False): ["verdana.ttf", "Verdana.ttf"],
        (True, False):  ["verdanab.ttf", "Verdana Bold.ttf", "Verdana-Bold.ttf"],
        (False, True):  ["verdanai.ttf", "Verdana Italic.ttf"],
        (True, True):   ["verdanaz.ttf", "Verdana Bold Italic.ttf"],
    },
    "Georgia": {
        (False, False): ["georgia.ttf", "Georgia.ttf"],
        (True, False):  ["georgiab.ttf", "Georgia Bold.ttf", "Georgia-Bold.ttf"],
        (False, True):  ["georgiai.ttf", "Georgia Italic.ttf"],
        (True, True):   ["georgiaz.ttf", "Georgia Bold Italic.ttf"],
    },
    "Tahoma": {
        (False, False): ["tahoma.ttf", "Tahoma.ttf"],
        (True, False):  ["tahomabd.ttf", "Tahoma Bold.ttf", "Tahoma-Bold.ttf"],
        (False, True):  ["tahoma.ttf", "Tahoma.ttf"],
        (True, True):   ["tahomabd.ttf", "Tahoma Bold.ttf"],
    },
    "Trebuchet MS": {
        (False, False): ["trebuc.ttf", "Trebuchet MS.ttf", "TrebuchetMS.ttf"],
        (True, False):  ["trebucbd.ttf", "Trebuchet MS Bold.ttf"],
        (False, True):  ["trebucit.ttf", "Trebuchet MS Italic.ttf"],
        (True, True):   ["trebucbi.ttf", "Trebuchet MS Bold Italic.ttf"],
    },
    "Cambria": {
        (False, False): ["cambria.ttc", "Cambria.ttf"],
        (True, False):  ["cambriab.ttf", "Cambria-Bold.ttf"],
        (False, True):  ["cambriai.ttf", "Cambria-Italic.ttf"],
        (True, True):   ["cambriaz.ttf", "Cambria-BoldItalic.ttf"],
    },
    "Candara": {
        (False, False): ["Candara.ttf", "candara.ttf"],
        (True, False):  ["Candarab.ttf", "Candara-Bold.ttf"],
        (False, True):  ["Candarai.ttf", "Candara-Italic.ttf"],
        (True, True):   ["Candaraz.ttf", "Candara-BoldItalic.ttf"],
    },
    "Corbel": {
        (False, False): ["corbel.ttf", "Corbel.ttf"],
        (True, False):  ["corbelb.ttf", "Corbel-Bold.ttf"],
        (False, True):  ["corbeli.ttf", "Corbel-Italic.ttf"],
        (True, True):   ["corbelz.ttf", "Corbel-BoldItalic.ttf"],
    },
    "Impact": {
        (False, False): ["impact.ttf", "Impact.ttf"],
        (True, False):  ["impact.ttf", "Impact.ttf"],
        (False, True):  ["impact.ttf", "Impact.ttf"],
        (True, True):   ["impact.ttf", "Impact.ttf"],
    },
    "Consolas": {
        (False, False): ["consola.ttf", "Consolas.ttf"],
        (True, False):  ["consolab.ttf", "Consolas-Bold.ttf"],
        (False, True):  ["consolai.ttf", "Consolas-Italic.ttf"],
        (True, True):   ["consolaz.ttf", "Consolas-BoldItalic.ttf"],
    },
    # Cross-platform fallback fonts (Linux/macOS)
    "DejaVu Sans": {
        (False, False): ["DejaVuSans.ttf"],
        (True, False):  ["DejaVuSans-Bold.ttf"],
        (False, True):  ["DejaVuSans-Oblique.ttf"],
        (True, True):   ["DejaVuSans-BoldOblique.ttf"],
    },
    "Liberation Sans": {
        (False, False): ["LiberationSans-Regular.ttf"],
        (True, False):  ["LiberationSans-Bold.ttf"],
        (False, True):  ["LiberationSans-Italic.ttf"],
        (True, True):   ["LiberationSans-BoldItalic.ttf"],
    },
    "Noto Sans": {
        (False, False): ["NotoSans-Regular.ttf", "NotoSans[wdth,wght].ttf"],
        (True, False):  ["NotoSans-Bold.ttf"],
        (False, True):  ["NotoSans-Italic.ttf"],
        (True, True):   ["NotoSans-BoldItalic.ttf"],
    },
}

# Cache: resolved font paths so we don't scan dirs repeatedly
_font_path_cache = {}


def _find_font_file(family, bold=False, italic=False):
    """Find a font file on disk for the given family + style. Returns path or None."""
    key = (family, bold, italic)
    if key in _font_path_cache:
        return _font_path_cache[key]

    style_key = (bool(bold), bool(italic))
    filenames = _FONT_FILENAMES.get(family, {}).get(style_key, [])
    # Also try regular variant as fallback
    if not filenames:
        filenames = _FONT_FILENAMES.get(family, {}).get((False, False), [])

    font_dirs = _get_font_dirs()
    for font_dir in font_dirs:
        for fname in filenames:
            # Check direct path
            fp = font_dir / fname
            if fp.exists():
                _font_path_cache[key] = str(fp)
                return str(fp)
            # Search subdirectories (Linux organizes fonts in subdirs)
            for sub in font_dir.rglob(fname):
                _font_path_cache[key] = str(sub)
                return str(sub)

    _font_path_cache[key] = None
    return None


# Fallback chain: try these families in order if requested family not found
_FALLBACK_FAMILIES = ["Arial", "Liberation Sans", "DejaVu Sans", "Noto Sans"]


def _get_bake_font(size, bold=False, italic=False, family="Arial"):
    """Get a TrueType font matching family + style, with cross-platform fallback."""
    # Try requested family first
    fp = _find_font_file(family, bold, italic)
    if fp:
        try:
            return ImageFont.truetype(fp, size)
        except (OSError, IOError):
            pass

    # Try fallback families
    for fb_family in _FALLBACK_FAMILIES:
        if fb_family == family:
            continue
        fp = _find_font_file(fb_family, bold, italic)
        if fp:
            try:
                return ImageFont.truetype(fp, size)
            except (OSError, IOError):
                continue

    # Try any regular font from fallbacks (ignore bold/italic)
    for fb_family in _FALLBACK_FAMILIES:
        fp = _find_font_file(fb_family, False, False)
        if fp:
            try:
                return ImageFont.truetype(fp, size)
            except (OSError, IOError):
                continue

    return ImageFont.load_default()


def _wrap_text_lines(text, font, max_w):
    """Word-wrap text to fit within max_w pixels. Preserves embedded \\n."""
    out = []
    for paragraph in text.split("\n"):
        words = paragraph.split(" ")
        line = ""
        for word in words:
            test = (line + " " + word).strip()
            bbox = font.getbbox(test)
            tw = (bbox[2] - bbox[0]) if bbox else 0
            if tw > max_w and line:
                out.append(line)
                line = word
            else:
                line = test
        out.append(line)  # always append, even if empty (preserves blank lines)
    return out


def _draw_chars_pillow(draw, text, x, y, font, fill, letter_spacing, scale):
    """Draw text character by character with letter spacing in Pillow."""
    cx = x
    ls_px = letter_spacing * scale
    for ch in text:
        draw.text((cx, y), ch, fill=fill, font=font)
        bbox = font.getbbox(ch)
        cw = (bbox[2] - bbox[0]) if bbox else 0
        cx += cw + ls_px


# ── Bake overlays into slide image ──────────────────────────────────────────

@app.route("/api/slide/<int:num>/bake", methods=["POST"])
def bake_overlays(num):
    """Burn overlays into the slide JPG image and clear them from data.
    This makes text edits permanent — the slide looks like the original
    but with changes baked in."""
    slide_files = _get_slide_files()
    if num < 1 or num > len(slide_files):
        return jsonify({"error": "Invalid slide"}), 400

    data = load_data()
    slide_data = data.get(str(num), {})
    ovs = slide_data.get("overlays", [])
    if not ovs:
        return jsonify({"ok": True, "msg": "No overlays to bake"})

    snapshot = _snapshot_before_destructive("bake")
    slide_path = slide_files[num - 1]
    img = Image.open(slide_path).convert("RGBA")
    w, h = img.size

    for ov in ovs:
        x = int(ov["x"] * w)
        y = int(ov["y"] * h)
        ow = int(ov["w"] * w)
        oh = int(ov["h"] * h)
        opacity = ov.get("opacity", 1)

        if ov["type"] == "rect":
            r, g, b = _parse_hex_color(ov.get("fillColor", "#2563EB"), (37, 99, 235))
            a = int(opacity * 255)
            layer = Image.new("RGBA", img.size, (0,0,0,0))
            ImageDraw.Draw(layer).rectangle([x, y, x+ow, y+oh], fill=(r,g,b,a))
            img = Image.alpha_composite(img, layer)

        elif ov["type"] == "text":
            draw = ImageDraw.Draw(img)
            bg = ov.get("bgColor", "transparent")
            if bg and bg != "transparent":
                cr, cg, cb = _parse_hex_color(bg)
                layer = Image.new("RGBA", img.size, (0,0,0,0))
                ImageDraw.Draw(layer).rectangle([x, y, x+ow, y+oh], fill=(cr,cg,cb,int(opacity*255)))
                img = Image.alpha_composite(img, layer)
                draw = ImageDraw.Draw(img)

            r, g, b = _parse_hex_color(ov.get("color", "#FFFFFF"), (255, 255, 255))
            raw_fs = ov.get("fontSize", 18)
            font_size = max(8, int(raw_fs * w / 933))
            is_bold = ov.get("bold", False)
            is_italic = ov.get("italic", False)
            font_family = ov.get("fontFamily", "Arial")
            font = _get_bake_font(font_size, is_bold, is_italic, font_family)

            text = ov.get("text", "")
            line_height_mul = ov.get("lineHeight", 1.3)
            letter_spacing = ov.get("letterSpacing", 0)
            text_transform = ov.get("textTransform", "none")
            list_style = ov.get("listStyle", "none")
            has_shadow = ov.get("shadow", False)
            shadow_color = _parse_hex_color(ov.get("shadowColor", "#000000"))
            has_outline = ov.get("outline", False)
            outline_color = _parse_hex_color(ov.get("outlineColor", "#000000"))
            outline_width = ov.get("outlineWidth", 1)
            has_underline = ov.get("underline", False)
            auto_fit = ov.get("autoFit", False)

            # Apply textTransform
            if text_transform == "uppercase":
                text = text.upper()
            elif text_transform == "lowercase":
                text = text.lower()
            elif text_transform == "capitalize":
                text = text.title()

            # Apply listStyle
            if list_style != "none":
                raw_lines = text.split("\n")
                new_lines = []
                for idx_l, ln in enumerate(raw_lines):
                    if list_style == "bullet":
                        new_lines.append("\u2022 " + ln)
                    elif list_style == "number":
                        new_lines.append(f"{idx_l+1}. " + ln)
                text = "\n".join(new_lines)

            lines = _wrap_text_lines(text, font, ow - 16)
            line_h = int(font_size * line_height_mul)

            # AutoFit: reduce font size until text fits
            if auto_fit:
                test_fs = font_size
                while test_fs > 8:
                    test_font = _get_bake_font(test_fs, is_bold, is_italic, font_family)
                    test_lines = _wrap_text_lines(text, test_font, ow - 16)
                    total_h = len(test_lines) * int(test_fs * line_height_mul) + 16
                    if total_h <= oh:
                        break
                    test_fs -= 1
                font_size = test_fs
                font = _get_bake_font(font_size, is_bold, is_italic, font_family)
                lines = _wrap_text_lines(text, font, ow - 16)
                line_h = int(font_size * line_height_mul)

            fill_color = (r, g, b, int(opacity*255))
            align = ov.get("align", "left")
            v_align = ov.get("verticalAlign", "top")
            total_text_h = max(0, len(lines) * line_h)
            if v_align == "center":
                ty = y + max(0, (oh - total_text_h) // 2)
            elif v_align == "bottom":
                ty = y + max(0, oh - total_text_h - 8)
            else:
                ty = y + 8

            sr, sg, sb = shadow_color if has_shadow else (0, 0, 0)
            olr, olg, olb = outline_color if has_outline else (0, 0, 0)

            for ln in lines:
                bbox = font.getbbox(ln)
                tw = (bbox[2] - bbox[0]) if bbox else 0
                if align == "center":
                    tx = x + (ow - tw) // 2
                elif align == "right":
                    tx = x + ow - tw - 8
                else:
                    tx = x + 8

                # Draw shadow (offset copy)
                if has_shadow:
                    if letter_spacing > 0:
                        _draw_chars_pillow(draw, ln, tx + 2, ty + 2, font, (sr, sg, sb, int(opacity*180)), letter_spacing, w / 933)
                    else:
                        draw.text((tx + 2, ty + 2), ln, fill=(sr, sg, sb, int(opacity*180)), font=font)

                # Draw outline (8-direction offset)
                if has_outline:
                    oc = (olr, olg, olb, int(opacity*255))
                    for odx in [-outline_width, 0, outline_width]:
                        for ody in [-outline_width, 0, outline_width]:
                            if odx == 0 and ody == 0:
                                continue
                            if letter_spacing > 0:
                                _draw_chars_pillow(draw, ln, tx + odx, ty + ody, font, oc, letter_spacing, w / 933)
                            else:
                                draw.text((tx + odx, ty + ody), ln, fill=oc, font=font)

                # Draw main text
                if letter_spacing > 0:
                    _draw_chars_pillow(draw, ln, tx, ty, font, fill_color, letter_spacing, w / 933)
                else:
                    draw.text((tx, ty), ln, fill=fill_color, font=font)

                # Draw underline
                if has_underline:
                    ul_y = ty + font_size + 2
                    draw.line([(tx, ul_y), (tx + tw, ul_y)], fill=fill_color, width=max(1, font_size // 14))

                ty += line_h

        elif ov["type"] == "image":
            src = ov.get("src", "")
            if src.startswith("data:"):
                b64_part = src.split(",", 1)[1] if "," in src else ""
                # Estimate decoded size before decoding (b64 = ~4/3 of decoded)
                if len(b64_part) * 3 // 4 > MAX_OVERLAY_IMG_BYTES:
                    continue
                try:
                    img_data = base64.b64decode(b64_part)
                except (ValueError, TypeError):
                    continue
                if len(img_data) > MAX_OVERLAY_IMG_BYTES:
                    continue
                overlay_img = Image.open(io.BytesIO(img_data)).convert("RGBA")
                overlay_img = overlay_img.resize((ow, oh), Image.BILINEAR)
                img.paste(overlay_img, (x, y), overlay_img)

    # Save back as RGB JPG
    img.convert("RGB").save(str(slide_path), "JPEG", quality=95)

    # Clear overlays from data (keep notes)
    data[str(num)] = {"overlays": [], "notes": slide_data.get("notes", "")}
    save_data(data)

    log_id = _log_op("bake", text=f"Bake {len(ovs)} overlays on slide {num}",
                     scope="current", slide_num=num, snapshot=snapshot)
    return jsonify({"ok": True, "snapshot": snapshot, "log_id": log_id})


# ── OCR Route ───────────────────────────────────────────────────────────────

@app.route("/api/ocr/<int:num>", methods=["POST"])
def ocr_slide(num):
    """Run EasyOCR on slide image, return detected text regions with bounding boxes."""
    if not HAS_OCR:
        return jsonify({"error": "EasyOCR is not installed. Run: pip install easyocr"}), 400

    slide_files = _get_slide_files()
    if num < 1 or num > len(slide_files):
        return jsonify({"error": "Invalid slide number"}), 400

    global _ocr_reader
    if _ocr_reader is None:
        _ocr_reader = easyocr.Reader(['en'], gpu=False, verbose=False)

    img_path = str(slide_files[num - 1])
    img = Image.open(img_path)
    w, h = img.size

    try:
        results = _ocr_reader.readtext(img_path)
    except Exception as e:
        return jsonify({"error": f"OCR failed: {e}"}), 500

    regions = []
    for (bbox, text, conf) in results:
        text = text.strip()
        if not text or conf < 0.3:
            continue
        # bbox is [[x1,y1],[x2,y2],[x3,y3],[x4,y4]] — take bounding rect
        xs = [p[0] for p in bbox]
        ys = [p[1] for p in bbox]
        x1, y1 = min(xs), min(ys)
        x2, y2 = max(xs), max(ys)
        regions.append({
            "text": text,
            "x": x1 / w,
            "y": y1 / h,
            "w": (x2 - x1) / w,
            "h": (y2 - y1) / h,
            "conf": round(conf * 100),
        })

    merged = _merge_ocr_regions(regions)
    return jsonify({"regions": merged, "raw_count": len(regions)})


def _merge_ocr_regions(regions):
    """Merge OCR word boxes on the same line into line-level regions."""
    if not regions:
        return []
    # Sort by Y then X
    regions.sort(key=lambda r: (round(r["y"], 2), r["x"]))
    merged = []
    current = None
    for r in regions:
        if current and abs(r["y"] - current["y"]) < 0.02 and r["x"] < current["x"] + current["w"] + 0.03:
            # Same line — extend
            new_right = max(current["x"] + current["w"], r["x"] + r["w"])
            current["w"] = new_right - current["x"]
            current["h"] = max(current["h"], r["h"])
            current["text"] += " " + r["text"]
        else:
            if current:
                merged.append(current)
            current = dict(r)
    if current:
        merged.append(current)
    return merged


# ── Sample background color route ───────────────────────────────────────────

@app.route("/api/sample-color/<int:num>", methods=["POST"])
def sample_color(num):
    """Sample BOTH background color and text color from a region.
    Background: sampled from a strip above the region.
    Text: darkest/lightest pixels inside the region (the ink)."""
    slide_files = _get_slide_files()
    if num < 1 or num > len(slide_files):
        return jsonify({"error": "Invalid slide number"}), 400

    payload = request.json
    img = Image.open(slide_files[num - 1]).convert("RGB")
    w, h = img.size

    x1 = int(payload["x"] * w)
    y1 = int(payload["y"] * h)
    x2 = int((payload["x"] + payload["w"]) * w)
    y2 = int((payload["y"] + payload["h"]) * h)
    x1, y1 = max(0, x1), max(0, y1)
    x2, y2 = min(w, x2), min(h, y2)

    # --- Background color: strip ABOVE the region ---
    pad = 5
    sy1 = max(0, y1 - pad - 3)
    sy2 = max(0, y1 - pad)
    if sy2 <= sy1:
        sy1, sy2 = max(0, y1 - 2), y1
    strip = img.crop((x1, sy1, x2, max(sy2, sy1 + 1)))
    bg_px = strip.resize((1, 1), Image.BILINEAR).getpixel((0, 0))
    bg_hex = f"#{bg_px[0]:02x}{bg_px[1]:02x}{bg_px[2]:02x}"

    # --- Text color: find the most different pixels from bg inside the region ---
    region = img.crop((x1, y1, x2, y2))
    rw, rh = region.size
    bg_lum = 0.299 * bg_px[0] + 0.587 * bg_px[1] + 0.114 * bg_px[2]

    # Sample pixels on a grid and find the one furthest from bg luminance
    best_dist = 0
    text_color = (0, 0, 0) if bg_lum > 128 else (255, 255, 255)
    step = max(1, min(rw, rh) // 20)
    for sx in range(0, rw, step):
        for sy in range(0, rh, step):
            px = region.getpixel((sx, sy))
            lum = 0.299 * px[0] + 0.587 * px[1] + 0.114 * px[2]
            dist = abs(lum - bg_lum)
            if dist > best_dist:
                best_dist = dist
                text_color = px

    txt_hex = f"#{text_color[0]:02x}{text_color[1]:02x}{text_color[2]:02x}"

    # --- Detect font weight via ink density ---
    total_px = 0
    ink_px = 0
    for col in range(0, rw, max(1, rw // 60)):
        for row in range(rh):
            px = region.getpixel((col, row))
            lum = 0.299 * px[0] + 0.587 * px[1] + 0.114 * px[2]
            total_px += 1
            if abs(lum - bg_lum) > 80:
                ink_px += 1

    density = ink_px / max(1, total_px)
    # Bold text: >18% ink density. Normal: <14%
    if density > 0.22:
        font_weight = "extrabold"
    elif density > 0.17:
        font_weight = "bold"
    elif density > 0.13:
        font_weight = "semibold"
    else:
        font_weight = "normal"

    # --- Alignment: horizontal centre of mass of ink pixels in the region ---
    # Compute mean x of all "ink" pixels (luminance distance from bg > threshold)
    align = "left"
    cap_height_px = max(8, int(rh * 0.7))
    is_italic = False
    try:
        import numpy as np
        arr = np.array(region, dtype=np.int16)
        lum = 0.299 * arr[..., 0] + 0.587 * arr[..., 1] + 0.114 * arr[..., 2]
        ink_mask = np.abs(lum - bg_lum) > 60
        if ink_mask.sum() > 0:
            # Horizontal centre of mass — normalise to 0..1 across the region width.
            xs = np.where(ink_mask)[1]
            ys = np.where(ink_mask)[0]
            cx_norm = xs.mean() / max(1, rw)
            if cx_norm < 0.4:
                align = "left"
            elif cx_norm > 0.6:
                align = "right"
            else:
                align = "center"
            # Tight cap height — vertical extent of ink, minus internal whitespace.
            ink_y_min, ink_y_max = int(ys.min()), int(ys.max())
            tight_h = max(1, ink_y_max - ink_y_min + 1)
            # Cap height is ~70% of glyph bounding box for typical fonts; the ink
            # vertical extent already excludes top/bottom whitespace, so use it directly.
            cap_height_px = max(8, tight_h)
            # Italic detection: compare horizontal position of ink at top vs bottom
            # of the ink band. Italic glyphs tilt right → top-row mean x > bottom-row mean x.
            if tight_h >= 6:
                top_band = ink_mask[ink_y_min:ink_y_min + tight_h // 3]
                bot_band = ink_mask[ink_y_max - tight_h // 3:ink_y_max + 1]
                if top_band.sum() > 0 and bot_band.sum() > 0:
                    top_x = np.where(top_band)[1].mean()
                    bot_x = np.where(bot_band)[1].mean()
                    # >2px shift normalised by cap height = likely italic
                    if (top_x - bot_x) > max(2, tight_h * 0.15):
                        is_italic = True
    except Exception:
        pass

    return jsonify({
        "color": bg_hex,
        "textColor": txt_hex,
        "fontWeight": font_weight,
        "align": align,
        "cap_height_px": cap_height_px,
        "is_italic": is_italic,
    })


# ── Image upload route ──────────────────────────────────────────────────────

@app.route("/api/upload-image", methods=["POST"])
def upload_image():
    """Accept an image file, return as base64 data URL."""
    if "file" not in request.files:
        return jsonify({"error": "No file"}), 400
    f = request.files["file"]
    try:
        img = Image.open(f.stream)
    except Exception:
        return jsonify({"error": "Invalid image file"}), 400
    if img.format and img.format not in ALLOWED_IMAGE_FORMATS:
        return jsonify({"error": f"Unsupported image format: {img.format}"}), 400
    img = img.convert("RGBA")
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    b64 = base64.b64encode(buf.read()).decode()
    return jsonify({"src": f"data:image/png;base64,{b64}", "w": img.width, "h": img.height})


@app.route("/api/logo/add-overlay", methods=["POST"])
def add_logo_overlay():
    """Add a logo as a draggable image overlay to one or all slides."""
    payload = request.json or {}
    src = payload.get("src")
    if not src:
        return jsonify({"error": "No src provided"}), 400
    try:
        x       = float(payload.get("x", 0))
        y       = float(payload.get("y", 0))
        w       = float(payload.get("w", 0.2))
        h       = float(payload.get("h", 0.2))
        opacity = float(payload.get("opacity", 1.0))
    except (TypeError, ValueError):
        return jsonify({"error": "Invalid overlay params"}), 400

    scope     = payload.get("scope", "all")
    slide_num = int(payload.get("slide_num", 1))

    data   = load_data()
    slides = _get_slide_files()

    targets = [str(slide_num)] if scope == "current" \
              else [str(i + 1) for i in range(len(slides))]

    ov = {"type": "image", "x": x, "y": y, "w": w, "h": h,
          "src": src, "opacity": opacity}
    for t in targets:
        data.setdefault(t, {"overlays": [], "notes": ""}) \
            .setdefault("overlays", []).append(ov)

    save_data(data)
    return jsonify({"ok": True, "count": len(targets)})


@app.route("/api/remove-background", methods=["POST"])
def remove_background():
    """Strip background from a base64 image overlay using rembg."""
    if not HAS_REMBG:
        return jsonify({"error": "rembg not installed — run: pip install rembg"}), 501
    payload = _ensure_dict(request.json)
    src = payload.get("src", "")
    if not src.startswith("data:image/"):
        return jsonify({"error": "Invalid image data"}), 400
    try:
        header, b64data = src.split(",", 1)
        raw = base64.b64decode(b64data)
        if len(raw) > MAX_OVERLAY_IMG_BYTES:
            return jsonify({"error": "Image too large (max 8 MB)"}), 413
        result = rembg_remove(raw)
        b64out = base64.b64encode(result).decode()
        return jsonify({"src": f"data:image/png;base64,{b64out}"})
    except Exception as e:
        return jsonify({"error": f"Background removal failed: {e}"}), 500


# ── Slide reorder route ─────────────────────────────────────────────────────

@app.route("/api/reorder", methods=["POST"])
def reorder_slides():
    """Reorder slide files on disk. Expects {"order": [3,1,2,...]}."""
    payload = _ensure_dict(request.json)
    new_order = payload.get("order", [])
    slide_files = _get_slide_files()

    if sorted(new_order) != list(range(1, len(slide_files) + 1)):
        return jsonify({"error": "Invalid order"}), 400

    snapshot = _snapshot_before_destructive("reorder")
    tmp_dir = SLIDES_DIR / "_reorder_tmp"
    try:
        if tmp_dir.exists():
            shutil.rmtree(str(tmp_dir))
        tmp_dir.mkdir()

        # Copy to temp with new names
        for new_idx, old_idx in enumerate(new_order, 1):
            src = SLIDES_DIR / f"slide-{old_idx:02d}.jpg"
            dst = tmp_dir / f"slide-{new_idx:02d}.jpg"
            if not src.exists():
                raise RuntimeError(f"Missing slide-{old_idx:02d}.jpg")
            shutil.copy2(str(src), str(dst))

        # All copies succeeded — swap
        for f in SLIDES_DIR.glob("slide-*.jpg"):
            f.unlink()
        for f in tmp_dir.glob("slide-*.jpg"):
            shutil.move(str(f), str(SLIDES_DIR / f.name))
    except Exception as e:
        return jsonify({"error": f"Reorder failed: {e}"}), 500
    finally:
        shutil.rmtree(str(tmp_dir), ignore_errors=True)

    # Reorder saved data too
    data = load_data()
    new_data = {}
    for new_idx, old_idx in enumerate(new_order, 1):
        new_data[str(new_idx)] = data.get(str(old_idx), {"overlays": [], "notes": ""})
    save_data(new_data)

    # Reorder comments in lock-step with slides
    with _data_lock:
        if COMMENTS_FILE.exists():
            try:
                cdata = json.loads(COMMENTS_FILE.read_text())
                new_comments = {}
                for new_idx, old_idx in enumerate(new_order, 1):
                    if str(old_idx) in cdata:
                        new_comments[str(new_idx)] = cdata[str(old_idx)]
                COMMENTS_FILE.write_text(json.dumps(new_comments, indent=2))
            except (json.JSONDecodeError, OSError):
                pass

    log_id = _log_op("reorder", text=f"Reordered {len(new_order)} slides",
                     scope="all", count=len(new_order), snapshot=snapshot)
    return jsonify({"ok": True, "snapshot": snapshot, "log_id": log_id})


# ── Deck info / slide structure ops ─────────────────────────────────────────

@app.route("/api/deck/info", methods=["GET"])
def deck_info():
    return jsonify({
        "deck_name": _get_deck_name(),
        "num_slides": len(_get_slide_files()),
    })


@app.route("/api/save", methods=["POST"])
def manual_save_checkpoint():
    """Drop a manual snapshot into the history chain."""
    snapshot = _snapshot_before_destructive("save")
    log_id = _log_op("save", text="Manual save checkpoint",
                     scope="all", count=len(_get_slide_files()), snapshot=snapshot)
    return jsonify({"ok": True, "snapshot": snapshot, "log_id": log_id})


@app.route("/api/slide/<int:num>/delete", methods=["POST"])
def delete_slide(num):
    """Delete a single slide, renumbering subsequent slides and rekeying
    overlay data + comments to stay aligned."""
    import copy as _copy
    slide_files = _get_slide_files()
    total = len(slide_files)
    if num < 1 or num > total:
        return jsonify({"error": "Invalid slide"}), 400
    if total <= 1:
        return jsonify({"error": "Cannot delete the only slide"}), 400

    snapshot = _snapshot_before_destructive("delete")
    try:
        target = SLIDES_DIR / f"slide-{num:02d}.jpg"
        if target.exists():
            target.unlink()
        # Shift subsequent files down by one
        for i in range(num + 1, total + 1):
            src = SLIDES_DIR / f"slide-{i:02d}.jpg"
            dst = SLIDES_DIR / f"slide-{i-1:02d}.jpg"
            if src.exists():
                shutil.move(str(src), str(dst))

        # Rekey slide_data.json
        with _data_lock:
            data = json.loads(DATA_FILE.read_text()) if DATA_FILE.exists() else {}
            new_data = {}
            for k, v in data.items():
                try:
                    idx = int(k)
                except (TypeError, ValueError):
                    continue
                if idx < num:
                    new_data[str(idx)] = v
                elif idx == num:
                    continue
                else:
                    new_data[str(idx - 1)] = v
            DATA_FILE.write_text(json.dumps(new_data, indent=2))

        # Rekey comments.json
        with _data_lock:
            if COMMENTS_FILE.exists():
                try:
                    cdata = json.loads(COMMENTS_FILE.read_text())
                    new_c = {}
                    for k, v in cdata.items():
                        try:
                            idx = int(k)
                        except (TypeError, ValueError):
                            continue
                        if idx < num:
                            new_c[str(idx)] = v
                        elif idx == num:
                            continue
                        else:
                            new_c[str(idx - 1)] = v
                    COMMENTS_FILE.write_text(json.dumps(new_c, indent=2))
                except (json.JSONDecodeError, OSError):
                    pass
    except Exception as e:
        return jsonify({"error": f"Delete failed: {e}"}), 500

    log_id = _log_op("delete-slide",
                     text=f"Deleted slide {num}",
                     scope="all", count=1, snapshot=snapshot)
    return jsonify({"ok": True, "num_slides": len(_get_slide_files()),
                    "snapshot": snapshot, "log_id": log_id})


@app.route("/api/slide/<int:num>/duplicate", methods=["POST"])
def duplicate_slide(num):
    """Duplicate a slide, inserting the clone at num+1 and shifting subsequent
    slides up by one. Overlay data and comments are deep-copied to the new index."""
    import copy as _copy
    slide_files = _get_slide_files()
    total = len(slide_files)
    if num < 1 or num > total:
        return jsonify({"error": "Invalid slide"}), 400

    snapshot = _snapshot_before_destructive("duplicate")
    try:
        # Shift slides num+1..total up by one
        for i in range(total, num, -1):
            src = SLIDES_DIR / f"slide-{i:02d}.jpg"
            dst = SLIDES_DIR / f"slide-{i+1:02d}.jpg"
            if src.exists():
                shutil.move(str(src), str(dst))
        # Copy the target file to num+1
        src = SLIDES_DIR / f"slide-{num:02d}.jpg"
        dst = SLIDES_DIR / f"slide-{num+1:02d}.jpg"
        shutil.copy2(str(src), str(dst))

        # Rekey slide_data.json
        with _data_lock:
            data = json.loads(DATA_FILE.read_text()) if DATA_FILE.exists() else {}
            new_data = {}
            for k, v in data.items():
                try:
                    idx = int(k)
                except (TypeError, ValueError):
                    continue
                if idx <= num:
                    new_data[str(idx)] = v
                else:
                    new_data[str(idx + 1)] = v
            if str(num) in data:
                new_data[str(num + 1)] = _copy.deepcopy(data[str(num)])
            else:
                new_data[str(num + 1)] = {"overlays": [], "notes": ""}
            DATA_FILE.write_text(json.dumps(new_data, indent=2))

        # Rekey comments.json
        with _data_lock:
            if COMMENTS_FILE.exists():
                try:
                    cdata = json.loads(COMMENTS_FILE.read_text())
                    new_c = {}
                    for k, v in cdata.items():
                        try:
                            idx = int(k)
                        except (TypeError, ValueError):
                            continue
                        if idx <= num:
                            new_c[str(idx)] = v
                        else:
                            new_c[str(idx + 1)] = v
                    if str(num) in cdata:
                        new_c[str(num + 1)] = _copy.deepcopy(cdata[str(num)])
                    COMMENTS_FILE.write_text(json.dumps(new_c, indent=2))
                except (json.JSONDecodeError, OSError):
                    pass
    except Exception as e:
        return jsonify({"error": f"Duplicate failed: {e}"}), 500

    log_id = _log_op("duplicate-slide",
                     text=f"Duplicated slide {num}",
                     scope="all", count=1, snapshot=snapshot)
    return jsonify({"ok": True, "num_slides": len(_get_slide_files()),
                    "new_slide": num + 1,
                    "snapshot": snapshot, "log_id": log_id})


@app.route("/api/slide/<int:num>/download.png", methods=["GET"])
def download_slide_png(num):
    """Convert the slide JPG to PNG and serve as an attachment."""
    slide_files = _get_slide_files()
    if num < 1 or num > len(slide_files):
        return jsonify({"error": "Invalid slide"}), 400
    src = SLIDES_DIR / f"slide-{num:02d}.jpg"
    if not src.exists():
        return jsonify({"error": "Slide not found"}), 404
    buf = io.BytesIO()
    Image.open(src).convert("RGB").save(buf, "PNG")
    buf.seek(0)
    deck = _get_deck_name() or "slide"
    base = re.sub(r"[^A-Za-z0-9_\-]", "_", Path(deck).stem)[:60] or "slide"
    return send_file(buf, mimetype="image/png", as_attachment=True,
                     download_name=f"{base}-slide-{num:02d}.png")


# ── PDF export route ────────────────────────────────────────────────────────

@app.route("/api/export-pdf", methods=["POST"])
def export_pdf():
    """Export slides as a multi-page PDF."""
    _cleanup_old_exports()
    slide_files = _get_slide_files()
    if not slide_files:
        return jsonify({"error": "No slides to export"}), 400

    first_img = Image.open(slide_files[0]).convert("RGB")
    rest = (Image.open(sf).convert("RGB") for sf in slide_files[1:])

    out_path = EXPORT_DIR / f"Slides_Export_{uuid.uuid4().hex[:8]}.pdf"
    first_img.save(str(out_path), save_all=True, append_images=rest, resolution=150)

    return send_file(str(out_path), as_attachment=True, download_name="Slides_Export.pdf")


# ── Image overlay in PPTX export ───────────────────────────────────────────

def _add_overlay(slide, ov, slide_w, slide_h):
    """Add a text, rect, arrow, or image overlay to a slide."""
    kind = ov.get("type", "text")

    left   = int(ov["x"] * slide_w)
    top    = int(ov["y"] * slide_h)
    width  = int(ov["w"] * slide_w)
    height = int(ov["h"] * slide_h)

    if kind == "text":
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = ov.get("text", "")

        font = run.font
        font.size = Pt(ov.get("fontSize", 18))
        r, g, b = _parse_hex_color(ov.get("color", "#FFFFFF"), (255, 255, 255))
        font.color.rgb = RGBColor(r, g, b)
        font.bold   = ov.get("bold", False)
        font.italic = ov.get("italic", False)

        align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
        p.alignment = align_map.get(ov.get("align","left"), PP_ALIGN.LEFT)

        bg_hex = ov.get("bgColor", "")
        if bg_hex and bg_hex != "transparent":
            br, bg_, bb = _parse_hex_color(bg_hex)
            fill = txBox.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(br, bg_, bb)

    elif kind == "rect":
        shape = slide.shapes.add_shape(1, left, top, width, height)
        r, g, b = _parse_hex_color(ov.get("fillColor", "#2563EB"), (37, 99, 235))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(r, g, b)
        shape.line.fill.background()

    elif kind == "image":
        src = ov.get("src", "")
        if src.startswith("data:"):
            b64_part = src.split(",", 1)[1] if "," in src else ""
            if len(b64_part) * 3 // 4 > MAX_OVERLAY_IMG_BYTES:
                return
            try:
                img_data = base64.b64decode(b64_part)
            except (ValueError, TypeError):
                return
            if len(img_data) > MAX_OVERLAY_IMG_BYTES:
                return
            img_stream = io.BytesIO(img_data)
            slide.shapes.add_picture(img_stream, left, top, width, height)

    elif kind == "arrow":
        line_shape = slide.shapes.add_connector(1, left, top, left+width, top+height)
        r, g, b = _parse_hex_color(ov.get("color", "#FF0000"), (255, 0, 0))
        line_shape.line.color.rgb = RGBColor(r, g, b)
        line_shape.line.width = Pt(2)


# ── Export as PNG ZIP ───────────────────────────────────────────────────────

@app.route("/api/export-png-zip", methods=["POST"])
def export_png_zip():
    """Export all slides as a ZIP of PNG images."""
    _cleanup_old_exports()
    import zipfile
    slide_files = _get_slide_files()
    if not slide_files:
        return jsonify({"error": "No slides"}), 400

    zip_path = EXPORT_DIR / f"slides_export_{uuid.uuid4().hex[:8]}.zip"
    with zipfile.ZipFile(str(zip_path), 'w', zipfile.ZIP_DEFLATED) as zf:
        for sf in slide_files:
            img = Image.open(sf).convert("RGB")
            png_buf = io.BytesIO()
            img.save(png_buf, format="PNG")
            png_buf.seek(0)
            zf.writestr(sf.stem + ".png", png_buf.read())

    return send_file(str(zip_path), as_attachment=True, download_name="slides_export.zip")


# ── Export as GIF ───────────────────────────────────────────────────────────

@app.route("/api/export-gif", methods=["POST"])
def export_gif():
    """Export slides as an animated GIF slideshow."""
    _cleanup_old_exports()
    payload = request.json or {}
    duration_ms = int(payload.get("duration", 2000))  # ms per slide

    slide_files = _get_slide_files()
    if not slide_files:
        return jsonify({"error": "No slides"}), 400

    frames = []
    for sf in slide_files:
        img = Image.open(sf).convert("RGB")
        img = img.resize((800, 450), Image.BILINEAR)
        frames.append(img)

    gif_path = EXPORT_DIR / f"slides_export_{uuid.uuid4().hex[:8]}.gif"
    frames[0].save(str(gif_path), save_all=True, append_images=frames[1:],
                   duration=duration_ms, loop=0, optimize=True)

    return send_file(str(gif_path), as_attachment=True, download_name="slides_export.gif")


# ── Image Filters ───────────────────────────────────────────────────────────

@app.route("/api/slide/<int:num>/filter", methods=["POST"])
def apply_filter(num):
    """Apply image filter to a slide. Filters: brightness, contrast, blur, grayscale, sepia."""
    from PIL import ImageEnhance, ImageFilter as PilFilter

    slide_files = _get_slide_files()
    if num < 1 or num > len(slide_files):
        return jsonify({"error": "Invalid slide"}), 400

    payload = _ensure_dict(request.json)
    filter_type = payload.get("filter", "")
    try:
        value = float(payload.get("value", 1.0))
    except (TypeError, ValueError):
        return jsonify({"error": "value must be numeric"}), 400

    img = Image.open(slide_files[num - 1]).convert("RGB")

    if filter_type == "brightness":
        img = ImageEnhance.Brightness(img).enhance(value)
    elif filter_type == "contrast":
        img = ImageEnhance.Contrast(img).enhance(value)
    elif filter_type == "saturation":
        img = ImageEnhance.Color(img).enhance(value)
    elif filter_type == "blur":
        img = img.filter(PilFilter.GaussianBlur(radius=value))
    elif filter_type == "sharpen":
        img = ImageEnhance.Sharpness(img).enhance(value)
    elif filter_type == "grayscale":
        import cv2, numpy as np
        arr = np.array(img)
        gray = cv2.cvtColor(arr, cv2.COLOR_RGB2GRAY)
        img = Image.fromarray(cv2.cvtColor(gray, cv2.COLOR_GRAY2RGB))
    elif filter_type == "sepia":
        import numpy as np
        arr = np.array(img, dtype=np.float64)
        sepia_filter = np.array([[0.393, 0.769, 0.189],
                                  [0.349, 0.686, 0.168],
                                  [0.272, 0.534, 0.131]])
        sepia = arr @ sepia_filter.T
        sepia = np.clip(sepia, 0, 255).astype(np.uint8)
        img = Image.fromarray(sepia)
    else:
        return jsonify({"error": f"Unknown filter: {filter_type}"}), 400

    img.save(str(slide_files[num - 1]), quality=95)
    return jsonify({"ok": True})


# ── Multi-filter pipeline ───────────────────────────────────────────────────

def _apply_filter_chain(img, params):
    """Apply a chain of filters to a PIL RGB image and return the result.
    params: dict with optional keys:
      brightness (1.0=neutral, 0..3),
      contrast   (1.0=neutral, 0..3),
      saturation (1.0=neutral, 0..3),
      hue        (degrees, -180..180),
      blur       (gaussian radius px, 0..30),
      sharpen    (0..3 = enhance factor),
      sepia      (0..1, intensity),
      grayscale  (0..1, intensity),
    Filters are applied in a stable order so behaviour matches the
    browser's live CSS preview.
    """
    from PIL import ImageEnhance, ImageFilter as PilFilter
    import numpy as np

    def f(key, default=0.0):
        try:
            return float(params.get(key, default))
        except (TypeError, ValueError):
            return default

    hue        = max(-180.0, min(180.0, f("hue", 0.0)))
    saturation = max(0.0,    min(3.0,   f("saturation", 1.0)))
    brightness = max(0.0,    min(3.0,   f("brightness", 1.0)))
    contrast   = max(0.0,    min(3.0,   f("contrast", 1.0)))
    sepia      = max(0.0,    min(1.0,   f("sepia", 0.0)))
    grayscale  = max(0.0,    min(1.0,   f("grayscale", 0.0)))
    blur       = max(0.0,    min(30.0,  f("blur", 0.0)))
    sharpen    = max(0.0,    min(3.0,   f("sharpen", 1.0)))

    # 1. Hue rotation (HSV)
    if abs(hue) > 0.5:
        hsv = np.array(img.convert("HSV"), dtype=np.int16)
        hsv[..., 0] = (hsv[..., 0] + int(hue * 255 / 360)) % 256
        img = Image.fromarray(hsv.astype(np.uint8), "HSV").convert("RGB")

    # 2. Saturation
    if abs(saturation - 1.0) > 1e-3:
        img = ImageEnhance.Color(img).enhance(saturation)

    # 3. Brightness
    if abs(brightness - 1.0) > 1e-3:
        img = ImageEnhance.Brightness(img).enhance(brightness)

    # 4. Contrast
    if abs(contrast - 1.0) > 1e-3:
        img = ImageEnhance.Contrast(img).enhance(contrast)

    # 5. Sepia (mix with sepia-toned version by intensity)
    if sepia > 0:
        arr = np.array(img, dtype=np.float64)
        m = np.array([[0.393, 0.769, 0.189],
                      [0.349, 0.686, 0.168],
                      [0.272, 0.534, 0.131]])
        toned = np.clip(arr @ m.T, 0, 255)
        mixed = arr * (1 - sepia) + toned * sepia
        img = Image.fromarray(np.clip(mixed, 0, 255).astype(np.uint8))

    # 6. Grayscale (mix luminance back into RGB by intensity)
    if grayscale > 0:
        arr = np.array(img, dtype=np.float64)
        lum = (0.299 * arr[..., 0] + 0.587 * arr[..., 1] + 0.114 * arr[..., 2])
        gray = np.stack([lum, lum, lum], axis=-1)
        mixed = arr * (1 - grayscale) + gray * grayscale
        img = Image.fromarray(np.clip(mixed, 0, 255).astype(np.uint8))

    # 7. Blur (Gaussian)
    if blur > 0.05:
        img = img.filter(PilFilter.GaussianBlur(radius=blur))

    # 8. Sharpen (Pillow's Sharpness enhancer, >1 sharpens, <1 softens)
    if abs(sharpen - 1.0) > 1e-3:
        img = ImageEnhance.Sharpness(img).enhance(sharpen)

    return img


@app.route("/api/slide/<int:num>/filters", methods=["POST"])
def apply_filters_chain(num):
    """Apply a chain of filters in one pass.

    Body keys: brightness, contrast, saturation, hue, blur, sharpen, sepia,
    grayscale (numeric, see _apply_filter_chain), plus:
      scope:         "current" (default) | "all"
      from_original: bool — if True (default), apply filters to the slide's
                     _originals/ copy (idempotent — re-applying with the same
                     values produces the same output, no cumulative blur).
                     If False, apply on top of the current live JPG.
    """
    slide_files = _get_slide_files()
    if num < 1 or num > len(slide_files):
        return jsonify({"error": "Invalid slide"}), 400

    payload = _ensure_dict(request.json)
    scope = str(payload.get("scope", "current"))
    from_original = bool(payload.get("from_original", True))

    if scope == "all":
        targets = list(range(1, len(slide_files) + 1))
    else:
        targets = [num]

    snapshot = _snapshot_before_destructive("filters")
    applied = 0
    for n in targets:
        live_path = slide_files[n - 1]
        orig_path = ORIGINALS_DIR / live_path.name
        source = orig_path if (from_original and orig_path.exists()) else live_path
        img = Image.open(source).convert("RGB")
        img = _apply_filter_chain(img, payload)
        img.save(str(live_path), "JPEG", quality=95)
        applied += 1

    # Log entry so it's revertable from the Applied tab
    entry_id = uuid.uuid4().hex[:12]
    label_params = {k: round(float(payload.get(k, 0)), 3)
                    for k in ("brightness", "contrast", "saturation", "hue",
                              "blur", "sharpen", "sepia", "grayscale")
                    if k in payload}
    _append_wm_log({
        "id": entry_id,
        "kind": "filters",
        "text": ", ".join(f"{k}:{v}" for k, v in label_params.items()) or "(no change)",
        "scope": "all" if scope == "all" else "current",
        "slide_num": num if scope != "all" else None,
        "count": applied,
        "snapshot": snapshot,
        "timestamp": datetime.datetime.now().isoformat(timespec="seconds"),
        "filter_params": label_params,
    })

    return jsonify({"ok": True, "count": applied, "snapshot": snapshot,
                    "log_id": entry_id, "source": "originals" if from_original else "live"})


# ── Crop / Rotate Slide ─────────────────────────────────────────────────────

@app.route("/api/slide/<int:num>/crop", methods=["POST"])
def crop_slide(num):
    """Crop a slide image. Expects normalized coords {x, y, w, h}."""
    slide_files = _get_slide_files()
    if num < 1 or num > len(slide_files):
        return jsonify({"error": "Invalid slide"}), 400

    payload = _ensure_dict(request.json)
    snapshot = _snapshot_before_destructive("crop")
    img = Image.open(slide_files[num - 1]).convert("RGB")
    w, h = img.size
    x1 = int(payload["x"] * w)
    y1 = int(payload["y"] * h)
    x2 = int((payload["x"] + payload["w"]) * w)
    y2 = int((payload["y"] + payload["h"]) * h)
    img = img.crop((x1, y1, x2, y2))
    # Resize back to standard dimensions
    img = img.resize((SLIDE_W_PX, SLIDE_H_PX), Image.LANCZOS)
    img.save(str(slide_files[num - 1]), quality=95)
    log_id = _log_op("crop", text=f"Crop slide {num}",
                     scope="current", slide_num=num, snapshot=snapshot)
    return jsonify({"ok": True, "snapshot": snapshot, "log_id": log_id})


@app.route("/api/slide/<int:num>/rotate", methods=["POST"])
def rotate_slide(num):
    """Rotate a slide image. Expects {angle: 90/180/270}."""
    slide_files = _get_slide_files()
    if num < 1 or num > len(slide_files):
        return jsonify({"error": "Invalid slide"}), 400

    payload = _ensure_dict(request.json)
    try:
        angle = int(payload.get("angle", 90))
    except (TypeError, ValueError):
        return jsonify({"error": "angle must be an integer"}), 400
    snapshot = _snapshot_before_destructive("rotate")
    img = Image.open(slide_files[num - 1]).convert("RGB")
    img = img.rotate(-angle, expand=True)
    img = img.resize((SLIDE_W_PX, SLIDE_H_PX), Image.LANCZOS)
    img.save(str(slide_files[num - 1]), quality=95)
    log_id = _log_op("rotate", text=f"Rotate slide {num} by {angle}°",
                     scope="current", slide_num=num, snapshot=snapshot)
    return jsonify({"ok": True, "snapshot": snapshot, "log_id": log_id})


# ── QR Code Generation ──────────────────────────────────────────────────────

@app.route("/api/qr-generate", methods=["POST"])
def generate_qr():
    """Generate a QR code as base64 PNG. Requires 'url' in payload."""
    payload = _ensure_dict(request.json)
    url = str(payload.get("url", ""))[:2000]
    if not url:
        return jsonify({"error": "No URL provided"}), 400

    try:
        import qrcode
        qr = qrcode.make(url)
        buf = io.BytesIO()
        qr.save(buf, format="PNG")
    except ImportError:
        # Fallback: generate simple QR-like placeholder
        img = Image.new("RGB", (200, 200), "white")
        draw = ImageDraw.Draw(img)
        draw.rectangle([10, 10, 190, 190], outline="black", width=3)
        draw.text((40, 90), "QR", fill="black")
        buf = io.BytesIO()
        img.save(buf, format="PNG")

    buf.seek(0)
    b64 = base64.b64encode(buf.read()).decode()
    return jsonify({"src": f"data:image/png;base64,{b64}"})


# ── Custom Watermark ────────────────────────────────────────────────────────

def _build_text_watermark_layer(w, h, *, text, color_rgb, opacity, position,
                                 font_scale, rotation_deg, tile_spacing):
    """Render a transparent RGBA layer with a text watermark.
    font_scale: fraction of slide width for the font height (e.g. 0.067 = w/15).
    rotation_deg: rotate the text. Applied to the whole layer for tiled/center,
                  to a sub-layer for corner positions.
    tile_spacing: multiplier for gap between tiled repeats (1.0 = default).
    """
    alpha = max(0, min(255, int(opacity * 255)))
    r, g, b = color_rgb
    fill = (r, g, b, alpha)
    font_size = max(8, int(w * font_scale))
    font = _get_bake_font(font_size, bold=True, family="Arial")
    bbox = font.getbbox(text)
    tw = bbox[2] - bbox[0]
    th = bbox[3] - bbox[1]

    layer = Image.new("RGBA", (w, h), (0, 0, 0, 0))

    def _draw_rotated_text(canvas, x, y):
        if rotation_deg % 360 == 0:
            ImageDraw.Draw(canvas).text((x, y), text, fill=fill, font=font)
            return
        # Draw text on its own transparent stamp, rotate, then paste.
        pad = max(tw, th) // 2 + 4
        stamp = Image.new("RGBA", (tw + pad * 2, th + pad * 2), (0, 0, 0, 0))
        ImageDraw.Draw(stamp).text((pad, pad), text, fill=fill, font=font)
        stamp = stamp.rotate(rotation_deg, resample=Image.BICUBIC, expand=True)
        sw, sh = stamp.size
        cx = x + tw // 2
        cy = y + th // 2
        canvas.alpha_composite(stamp, (cx - sw // 2, cy - sh // 2))

    if position == "tiled":
        step_x = max(1, int((tw + 100) * max(0.3, tile_spacing)))
        step_y = max(1, int((th + 80) * max(0.3, tile_spacing)))
        for tx in range(0, w + step_x, step_x):
            for ty in range(0, h + step_y, step_y):
                _draw_rotated_text(layer, tx, ty)
    elif position == "center":
        _draw_rotated_text(layer, (w - tw) // 2, (h - th) // 2)
    elif position == "bottom-right":
        _draw_rotated_text(layer, w - tw - 20, h - th - 20)
    elif position == "bottom-left":
        _draw_rotated_text(layer, 20, h - th - 20)
    elif position == "top-right":
        _draw_rotated_text(layer, w - tw - 20, 20)
    elif position == "top-left":
        _draw_rotated_text(layer, 20, 20)
    else:
        _draw_rotated_text(layer, (w - tw) // 2, (h - th) // 2)
    return layer


# ── Applied watermark log ───────────────────────────────────────────────────
WATERMARK_LOG = BASE_DIR / "watermarks_applied.json"
# Snapshots created before this epoch time are hidden from the Applied list.
# Bumped by POST /api/watermarks/clear-log so the user can wipe the visible
# list without us deleting the on-disk history snapshots.
CLEARED_ORPHANS_MARKER = BASE_DIR / "cleared_orphans_before.txt"


def _get_cleared_orphans_before():
    try:
        return float(CLEARED_ORPHANS_MARKER.read_text().strip())
    except (OSError, ValueError):
        return 0.0


def _set_cleared_orphans_before(t):
    try:
        CLEARED_ORPHANS_MARKER.write_text(str(t))
    except OSError:
        pass


def _load_wm_log():
    with _data_lock:
        if WATERMARK_LOG.exists():
            try:
                return json.loads(WATERMARK_LOG.read_text())
            except (json.JSONDecodeError, OSError):
                return []
        return []


def _save_wm_log(entries):
    with _data_lock:
        WATERMARK_LOG.write_text(json.dumps(entries, indent=2))


def _append_wm_log(entry):
    """Append an op entry. A new op invalidates any undone (redo-able) entries
    in front of it, since the timeline has just branched — those redo snapshots
    no longer reflect a reachable future state."""
    with _data_lock:
        entries = []
        if WATERMARK_LOG.exists():
            try:
                entries = json.loads(WATERMARK_LOG.read_text())
            except (json.JSONDecodeError, OSError):
                entries = []
        # Drop any undone-but-not-yet-redone entries (the future was rewritten).
        entries = [e for e in entries if not e.get("undone")]
        entries.append(entry)
        # Cap at 100 most recent
        if len(entries) > 100:
            entries = entries[-100:]
        WATERMARK_LOG.write_text(json.dumps(entries, indent=2))


def _restore_from_snapshot(version_dir: Path, *, scope: str, slide_num=None):
    """Restore slides + data.json from a snapshot directory, respecting scope.

    scope == "current" with slide_num → only that one slide.
    Anything else → full restore (all slides + data.json).

    Raises RuntimeError if a required slide is missing or copy fails, so the
    caller can surface a 5xx instead of silently leaving disk in a half-state.
    """
    data_file = version_dir / "data.json"
    if scope == "current" and slide_num:
        try:
            n = int(slide_num)
        except (TypeError, ValueError):
            raise RuntimeError(f"Invalid slide_num: {slide_num!r}")
        name = f"slide-{n:02d}.jpg"
        src = version_dir / name
        dst = SLIDES_DIR / name
        if not src.exists():
            raise RuntimeError(f"Snapshot missing {name}")
        try:
            shutil.copy2(str(src), str(dst))
        except OSError as e:
            raise RuntimeError(f"Copy failed for {name}: {e}") from e
        if data_file.exists():
            with _data_lock:
                live = json.loads(DATA_FILE.read_text()) if DATA_FILE.exists() else {}
                snap_data = json.loads(data_file.read_text())
                live[str(n)] = snap_data.get(
                    str(n), {"overlays": [], "notes": ""})
                DATA_FILE.write_text(json.dumps(live, indent=2))
    else:
        snap_slides = sorted(version_dir.glob("slide-*.jpg"))
        if not snap_slides:
            raise RuntimeError("Snapshot contains no slide images")
        for f in SLIDES_DIR.glob("slide-*.jpg"):
            try:
                f.unlink()
            except OSError:
                pass  # best-effort cleanup; subsequent copy will overwrite
        for sf in snap_slides:
            try:
                shutil.copy2(str(sf), str(SLIDES_DIR / sf.name))
            except OSError as e:
                raise RuntimeError(f"Copy failed for {sf.name}: {e}") from e
        if data_file.exists():
            with _data_lock:
                DATA_FILE.write_text(data_file.read_text())
        comments_snap = version_dir / "comments.json"
        if comments_snap.exists():
            with _data_lock:
                COMMENTS_FILE.write_text(comments_snap.read_text())


def _take_redo_snapshot():
    """Capture the current state so a subsequent redo can restore it.
    Reused for undo's redo-stack; tagged with reason='redo-snap' so it doesn't
    show up in the Applied tab orphan list."""
    return _snapshot_before_destructive("redo-snap")


def _log_op(kind, *, text, scope, slide_num=None, count=1, snapshot=None):
    """Append a generic destructive-operation entry to the watermark log so
    it surfaces in the Applied tab and can be reverted via /watermarks/revert.
    Returns the entry id (or None if snapshot is missing)."""
    if not snapshot:
        return None
    entry_id = uuid.uuid4().hex[:12]
    _append_wm_log({
        "id": entry_id,
        "kind": kind,
        "text": text,
        "scope": scope,
        "slide_num": slide_num if scope == "current" else None,
        "count": count,
        "snapshot": snapshot,
        "timestamp": datetime.datetime.now().isoformat(timespec="seconds"),
    })
    return entry_id


def _snapshot_before_destructive(reason="watermark"):
    """Save a quick history snapshot. Returns version id or None."""
    try:
        ts = datetime.datetime.now().strftime("%Y%m%d_%H%M%S") + f"_{uuid.uuid4().hex[:4]}"
        version_dir = HISTORY_DIR / ts
        version_dir.mkdir(exist_ok=True)
        for sf in _get_slide_files():
            shutil.copy2(str(sf), str(version_dir / sf.name))
        if DATA_FILE.exists():
            shutil.copy2(str(DATA_FILE), str(version_dir / "data.json"))
        if COMMENTS_FILE.exists():
            shutil.copy2(str(COMMENTS_FILE), str(version_dir / "comments.json"))
        (version_dir / "_reason.txt").write_text(reason)
        return ts
    except OSError:
        return None


@app.route("/api/watermark", methods=["POST"])
def add_watermark():
    """Add a text watermark. Accepts:
        text, opacity, position (center/bottom-right/bottom-left/top-right/top-left/tiled),
        color (#rrggbb), font_scale (0.02–0.2), rotation (-90..90 deg),
        tile_spacing (0.5..3.0), scope ("all"|"current"), slide_num (1-based if current).
    """
    payload = _ensure_dict(request.json)
    text = str(payload.get("text", "CONFIDENTIAL"))[:200]
    if not text.strip():
        return jsonify({"error": "Text cannot be empty"}), 400
    try:
        opacity = float(payload.get("opacity", 0.15))
    except (TypeError, ValueError):
        opacity = 0.15
    opacity = max(0.02, min(1.0, opacity))
    position = str(payload.get("position", "center"))
    color_rgb = _parse_hex_color(payload.get("color", "#808080"), (128, 128, 128))
    try:
        font_scale = float(payload.get("font_scale", 1 / 15))
    except (TypeError, ValueError):
        font_scale = 1 / 15
    font_scale = max(0.02, min(0.25, font_scale))
    try:
        rotation = float(payload.get("rotation", 0))
    except (TypeError, ValueError):
        rotation = 0
    rotation = max(-90.0, min(90.0, rotation))
    try:
        tile_spacing = float(payload.get("tile_spacing", 1.0))
    except (TypeError, ValueError):
        tile_spacing = 1.0
    tile_spacing = max(0.3, min(4.0, tile_spacing))

    scope = str(payload.get("scope", "all"))
    all_files = _get_slide_files()
    skip_set = _parse_skip_slides(payload.get("skip_slides"))
    if scope == "current":
        try:
            n = int(payload.get("slide_num", 1))
        except (TypeError, ValueError):
            n = 1
        if n < 1 or n > len(all_files):
            return jsonify({"error": "Invalid slide_num"}), 400
        if n in skip_set:
            return jsonify({"error": f"Slide {n} is in the skip list"}), 400
        targets = [all_files[n - 1]]
    else:
        targets = [f for i, f in enumerate(all_files, start=1) if i not in skip_set]
        if not targets:
            return jsonify({"error": "Skip list excludes every slide — nothing to do"}), 400

    snapshot = _snapshot_before_destructive("watermark-text")
    for sf in targets:
        img = Image.open(sf).convert("RGBA")
        w, h = img.size
        layer = _build_text_watermark_layer(
            w, h, text=text, color_rgb=color_rgb, opacity=opacity,
            position=position, font_scale=font_scale, rotation_deg=rotation,
            tile_spacing=tile_spacing,
        )
        img = Image.alpha_composite(img, layer)
        img.convert("RGB").save(str(sf), quality=95)

    entry_id = uuid.uuid4().hex[:12]
    _append_wm_log({
        "id": entry_id,
        "kind": "text",
        "text": text,
        "color": "#{:02x}{:02x}{:02x}".format(*color_rgb),
        "opacity": opacity,
        "position": position,
        "rotation": rotation,
        "font_scale": font_scale,
        "tile_spacing": tile_spacing,
        "scope": scope,
        "slide_num": int(payload.get("slide_num", 1)) if scope == "current" else None,
        "skip_slides": sorted(skip_set) if skip_set else [],
        "count": len(targets),
        "snapshot": snapshot,
        "timestamp": datetime.datetime.now().isoformat(timespec="seconds"),
    })
    return jsonify({"ok": True, "count": len(targets),
                    "skipped": sorted(skip_set) if skip_set else [],
                    "snapshot": snapshot, "log_id": entry_id})


@app.route("/api/watermark/preview/<int:num>", methods=["POST"])
def preview_watermark(num):
    """Render a non-destructive preview of the text watermark on slide <num>.
    Returns base64 JPEG."""
    slide_files = _get_slide_files()
    if num < 1 or num > len(slide_files):
        return jsonify({"error": "Invalid slide"}), 400

    payload = _ensure_dict(request.json)
    text = str(payload.get("text", "CONFIDENTIAL"))[:200]
    try:
        opacity = max(0.02, min(1.0, float(payload.get("opacity", 0.15))))
    except (TypeError, ValueError):
        opacity = 0.15
    position = str(payload.get("position", "center"))
    color_rgb = _parse_hex_color(payload.get("color", "#808080"), (128, 128, 128))
    try:
        font_scale = max(0.02, min(0.25, float(payload.get("font_scale", 1 / 15))))
    except (TypeError, ValueError):
        font_scale = 1 / 15
    try:
        rotation = max(-90.0, min(90.0, float(payload.get("rotation", 0))))
    except (TypeError, ValueError):
        rotation = 0
    try:
        tile_spacing = max(0.3, min(4.0, float(payload.get("tile_spacing", 1.0))))
    except (TypeError, ValueError):
        tile_spacing = 1.0

    # Build at a smaller size for speed
    full = Image.open(slide_files[num - 1]).convert("RGBA")
    fw, fh = full.size
    target_w = 800
    scale = target_w / fw
    pw, ph = target_w, int(fh * scale)
    small = full.resize((pw, ph), Image.BILINEAR)
    layer = _build_text_watermark_layer(
        pw, ph, text=text, color_rgb=color_rgb, opacity=opacity,
        position=position, font_scale=font_scale, rotation_deg=rotation,
        tile_spacing=tile_spacing,
    )
    composed = Image.alpha_composite(small, layer).convert("RGB")
    buf = io.BytesIO()
    composed.save(buf, format="JPEG", quality=80)
    buf.seek(0)
    b64 = base64.b64encode(buf.read()).decode()
    return jsonify({"preview": f"data:image/jpeg;base64,{b64}"})


# ── Watermark Detection ─────────────────────────────────────────────────────

@app.route("/api/detect-watermark/<int:num>", methods=["POST"])
def detect_watermark(num):
    """Detect watermark regions by cross-slide consistency.

    True watermarks repeat in the same position across multiple slides with
    near-identical pixels; regular text content varies. We sample up to 6 slides,
    score each corner by how similar that corner looks across the sample, and
    only flag a corner if the cross-slide similarity is high AND the corner has
    real ink (not just a uniform background).
    """
    import cv2
    import numpy as np

    slide_files = _get_slide_files()
    if num < 1 or num > len(slide_files):
        return jsonify({"error": "Invalid slide"}), 400

    # Sample up to 6 slides for cross-slide comparison
    sample_idxs = list(range(len(slide_files)))
    if len(sample_idxs) > 6:
        step = max(1, len(sample_idxs) // 6)
        sample_idxs = sample_idxs[::step][:6]
        if (num - 1) not in sample_idxs:
            sample_idxs[0] = num - 1
    sample_imgs = [cv2.imread(str(slide_files[i])) for i in sample_idxs]
    sample_imgs = [s for s in sample_imgs if s is not None]
    if not sample_imgs:
        return jsonify({"candidates": []})

    # Resize all to a common reference size for comparison
    ref_h, ref_w = sample_imgs[0].shape[:2]
    norm = [cv2.resize(s, (ref_w, ref_h)) if s.shape[:2] != (ref_h, ref_w) else s for s in sample_imgs]

    corners = {
        "bottom-right": (int(ref_w*0.7), int(ref_h*0.85), ref_w, ref_h),
        "bottom-left":  (0, int(ref_h*0.85), int(ref_w*0.3), ref_h),
        "top-right":    (int(ref_w*0.7), 0, ref_w, int(ref_h*0.15)),
        "top-left":     (0, 0, int(ref_w*0.3), int(ref_h*0.15)),
    }

    candidates = []
    MIN_CONFIDENCE = 60  # require strong cross-slide consistency

    for corner_name, (x1, y1, x2, y2) in corners.items():
        crops = [cv2.cvtColor(s[y1:y2, x1:x2], cv2.COLOR_BGR2GRAY) for s in norm]
        if len(crops) < 2:
            break
        # Pairwise mean abs diff against the first sample
        ref = crops[0].astype(np.float32)
        diffs = [float(np.mean(np.abs(c.astype(np.float32) - ref))) for c in crops[1:]]
        mean_diff = float(np.mean(diffs)) if diffs else 0.0
        # Ink check: corner must have some non-background content
        edges = cv2.Canny(crops[0], 50, 150)
        edge_density = float(np.sum(edges > 0)) / max(1, edges.size)

        # Convert mean_diff to a 0..100 similarity score; <5 px diff is very similar
        similarity = max(0.0, 100.0 - mean_diff * 10.0)
        # Require both similarity AND ink density in watermark range
        if similarity >= MIN_CONFIDENCE and 0.01 < edge_density < 0.35:
            # Tight contour bbox in the corner of the *current* slide
            gray_now = cv2.cvtColor(norm[0], cv2.COLOR_BGR2GRAY)
            roi = gray_now[y1:y2, x1:x2]
            ed = cv2.Canny(roi, 50, 150)
            contours, _ = cv2.findContours(ed, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
            if contours:
                all_pts = np.vstack(contours)
                bx, by, bw, bh = cv2.boundingRect(all_pts)
                abs_x = (x1 + bx) / ref_w
                abs_y = (y1 + by) / ref_h
                abs_w = bw / ref_w
                abs_h = bh / ref_h
                pad = 0.01
                abs_x = max(0, abs_x - pad)
                abs_y = max(0, abs_y - pad)
                abs_w = min(1 - abs_x, abs_w + pad * 2)
                abs_h = min(1 - abs_y, abs_h + pad * 2)
                candidates.append({
                    "location": corner_name,
                    "x": round(abs_x, 4),
                    "y": round(abs_y, 4),
                    "w": round(abs_w, 4),
                    "h": round(abs_h, 4),
                    "confidence": round(similarity, 1),
                    "note": f"matches across {len(crops)} slides",
                })

    candidates.sort(key=lambda c: c["confidence"], reverse=True)
    return jsonify({"candidates": candidates, "sampled_slides": len(sample_imgs)})


@app.route("/api/remove-watermark/<int:num>", methods=["POST"])
def remove_watermark(num):
    """Remove a specific watermark region from a slide using inpainting."""
    import cv2
    import numpy as np

    slide_files = _get_slide_files()
    if num < 1 or num > len(slide_files):
        return jsonify({"error": "Invalid slide"}), 400

    payload = request.json
    regions = payload.get("regions", [])
    if not regions:
        return jsonify({"error": "No regions specified"}), 400

    snapshot = _snapshot_before_destructive("watermark-remove")
    img = cv2.imread(str(slide_files[num - 1]))
    h, w = img.shape[:2]

    mask = np.zeros((h, w), dtype=np.uint8)
    for r in regions:
        x1 = int(r["x"] * w)
        y1 = int(r["y"] * h)
        x2 = int((r["x"] + r["w"]) * w)
        y2 = int((r["y"] + r["h"]) * h)
        mask[y1:y2, x1:x2] = 255

    result = cv2.inpaint(img, mask, inpaintRadius=5, flags=cv2.INPAINT_TELEA)
    cv2.imwrite(str(slide_files[num - 1]), result, [cv2.IMWRITE_JPEG_QUALITY, 95])

    return jsonify({"ok": True, "snapshot": snapshot})


@app.route("/api/remove-watermark-all", methods=["POST"])
def remove_watermark_all():
    """Remove watermark from ALL slides at the same region."""
    import cv2
    import numpy as np

    payload = request.json
    regions = payload.get("regions", [])
    if not regions:
        return jsonify({"error": "No regions specified"}), 400

    slide_files = _get_slide_files()
    snapshot = _snapshot_before_destructive("watermark-remove-all")
    count = 0
    for sf in slide_files:
        img = cv2.imread(str(sf))
        h, w = img.shape[:2]
        mask = np.zeros((h, w), dtype=np.uint8)
        for r in regions:
            x1 = int(r["x"] * w)
            y1 = int(r["y"] * h)
            x2 = int((r["x"] + r["w"]) * w)
            y2 = int((r["y"] + r["h"]) * h)
            mask[y1:y2, x1:x2] = 255
        result = cv2.inpaint(img, mask, inpaintRadius=5, flags=cv2.INPAINT_TELEA)
        cv2.imwrite(str(sf), result, [cv2.IMWRITE_JPEG_QUALITY, 95])
        count += 1

    return jsonify({"ok": True, "slides": count, "snapshot": snapshot})


# ── Custom Watermark (Image) ───────────────────────────────────────────────

@app.route("/api/watermark-image", methods=["POST"])
def add_image_watermark():
    """Add an image watermark (logo). Accepts scope=all|current and slide_num."""
    if "image" not in request.files:
        return jsonify({"error": "No image uploaded"}), 400

    try:
        wm_img = Image.open(request.files["image"].stream).convert("RGBA")
    except Exception:
        return jsonify({"error": "Invalid image"}), 400
    try:
        wm_opacity = max(0.02, min(1.0, float(request.form.get("opacity", 0.3))))
    except (TypeError, ValueError):
        wm_opacity = 0.3
    wm_position = request.form.get("position", "bottom-right")
    try:
        wm_scale = max(0.02, min(0.6, float(request.form.get("scale", 0.15))))
    except (TypeError, ValueError):
        wm_scale = 0.15
    try:
        wm_padding = max(0, min(200, int(request.form.get("padding", 20))))
    except (TypeError, ValueError):
        wm_padding = 20

    scope = request.form.get("scope", "all")
    all_files = _get_slide_files()
    skip_set = _parse_skip_slides(request.form.get("skip_slides"))
    if scope == "current":
        try:
            n = int(request.form.get("slide_num", 1))
        except (TypeError, ValueError):
            n = 1
        if n < 1 or n > len(all_files):
            return jsonify({"error": "Invalid slide_num"}), 400
        if n in skip_set:
            return jsonify({"error": f"Slide {n} is in the skip list"}), 400
        slide_files = [all_files[n - 1]]
    else:
        slide_files = [f for i, f in enumerate(all_files, start=1) if i not in skip_set]
        if not slide_files:
            return jsonify({"error": "Skip list excludes every slide — nothing to do"}), 400

    snapshot = _snapshot_before_destructive("watermark-image")
    for sf in slide_files:
        img = Image.open(sf).convert("RGBA")
        w, h = img.size

        # Scale watermark
        wm_w = int(w * wm_scale)
        wm_h = int(wm_w * wm_img.height / wm_img.width)
        wm_resized = wm_img.resize((wm_w, wm_h), Image.LANCZOS)

        # Apply opacity
        r, g, b, a = wm_resized.split()
        import numpy as np
        a_arr = np.array(a).astype(float) * wm_opacity
        a = Image.fromarray(a_arr.astype(np.uint8))
        wm_resized = Image.merge("RGBA", (r, g, b, a))

        # Position
        cx = (w - wm_w) // 2
        cy = (h - wm_h) // 2
        pad = wm_padding
        _positions = {
            "top-left":      (pad, pad),
            "top-center":    (cx, pad),
            "top-right":     (w - wm_w - pad, pad),
            "middle-left":   (pad, cy),
            "center":        (cx, cy),
            "middle-right":  (w - wm_w - pad, cy),
            "bottom-left":   (pad, h - wm_h - pad),
            "bottom-center": (cx, h - wm_h - pad),
            "bottom-right":  (w - wm_w - pad, h - wm_h - pad),
        }
        if wm_position == "tiled":
            overlay = Image.new("RGBA", (w, h), (0, 0, 0, 0))
            for tx in range(0, w, wm_w + 60):
                for ty in range(0, h, wm_h + 40):
                    overlay.paste(wm_resized, (tx, ty), wm_resized)
            img = Image.alpha_composite(img, overlay)
            img.convert("RGB").save(str(sf), quality=95)
            continue
        pos = _positions.get(wm_position, _positions["bottom-right"])

        img.paste(wm_resized, pos, wm_resized)
        img.convert("RGB").save(str(sf), quality=95)

    entry_id = uuid.uuid4().hex[:12]
    _append_wm_log({
        "id": entry_id,
        "kind": "image",
        "filename": request.files["image"].filename or "image",
        "opacity": wm_opacity,
        "position": wm_position,
        "scale": wm_scale,
        "scope": scope,
        "slide_num": int(request.form.get("slide_num", 1)) if scope == "current" else None,
        "skip_slides": sorted(skip_set) if skip_set else [],
        "count": len(slide_files),
        "snapshot": snapshot,
        "timestamp": datetime.datetime.now().isoformat(timespec="seconds"),
    })
    return jsonify({"ok": True, "count": len(slide_files), "snapshot": snapshot, "log_id": entry_id})


# ── Full reset (restore originals, clear all edits) ────────────────────────

@app.route("/api/reset-all", methods=["POST"])
def reset_all():
    """Restore every slide from its `_originals/` copy and clear all
    user-applied state: overlays, notes, comments, watermark log.

    If no `_originals/` exist (e.g. PPTX uploaded before this feature shipped),
    falls back to the OLDEST history snapshot — the next-best approximation of
    'the original'.

    Does NOT touch: history snapshots, saved templates, or uploads.
    A history snapshot is taken first so the action itself is reversible.
    """
    originals = sorted(ORIGINALS_DIR.glob("slide-*.jpg"))
    source_label = "originals"

    if not originals:
        # Fallback: oldest history snapshot
        snapshot_dirs = sorted(
            [d for d in HISTORY_DIR.iterdir() if d.is_dir()],
            key=lambda d: d.stat().st_mtime
        ) if HISTORY_DIR.exists() else []
        if not snapshot_dirs:
            return jsonify({
                "error": "No originals and no history snapshots — nothing to reset to. Re-upload the PPTX to enable Reset.",
            }), 400
        oldest = snapshot_dirs[0]
        originals = sorted(oldest.glob("slide-*.jpg"))
        if not originals:
            return jsonify({"error": "Oldest history snapshot is empty"}), 400
        source_label = f"oldest snapshot ({oldest.name})"

    snapshot = _snapshot_before_destructive("reset-all")

    # Wipe live slides, copy source back
    for f in SLIDES_DIR.glob("slide-*.jpg"):
        f.unlink()
    restored = 0
    for orig in originals:
        shutil.copy2(str(orig), str(SLIDES_DIR / orig.name))
        restored += 1

    # Clear per-slide overlays/notes, comments, watermark log — under the lock
    with _data_lock:
        if DATA_FILE.exists():
            DATA_FILE.unlink()
        if COMMENTS_FILE.exists():
            COMMENTS_FILE.unlink()
        if WATERMARK_LOG.exists():
            WATERMARK_LOG.unlink()

    return jsonify({
        "ok": True,
        "slides_restored": restored,
        "snapshot": snapshot,
        "source": source_label,
    })


# ── Applied watermark list / revert ────────────────────────────────────────

@app.route("/api/watermarks/applied", methods=["GET"])
def list_applied_watermarks():
    """List recent watermark apply operations the user can revert.

    Merges two sources:
      1. The watermarks_applied.json log (full metadata).
      2. History snapshots tagged with a watermark-* reason but not in the log
         (orphans — from a previous session, or after Clear Log).
    """
    entries = _load_wm_log()
    logged_by_snapshot = {e.get("snapshot"): e for e in entries if e.get("snapshot")}
    cleared_before = _get_cleared_orphans_before()

    # Scan history dirs for watermark-related snapshots
    orphans = []
    if HISTORY_DIR.exists():
        for vd in HISTORY_DIR.iterdir():
            if not vd.is_dir():
                continue
            if vd.name in logged_by_snapshot:
                continue
            reason_file = vd / "_reason.txt"
            if not reason_file.exists():
                continue
            try:
                reason = reason_file.read_text().strip()
            except OSError:
                continue
            # Show snapshots from any tracked destructive op (watermark, filters,
            # crop, rotate, bake, remove-logo, find-replace, reorder, load-template,
            # restore-version). Reset snapshots are intentionally hidden — they're
            # the snapshot taken by Reset itself, not a user op to revert.
            if reason in ("reset-all", "redo-snap") or not reason:
                continue
            tracked_prefixes = ("watermark", "filters", "crop", "rotate", "bake",
                                "remove-logo", "find-replace", "reorder",
                                "load-template", "restore-version", "inpaint-region")
            if not reason.startswith(tracked_prefixes):
                continue
            try:
                mtime = vd.stat().st_mtime
                ts = datetime.datetime.fromtimestamp(mtime).isoformat(timespec="seconds")
            except OSError:
                ts = ""
                mtime = 0
            # Skip orphans hidden by a prior "Clear log" action
            if mtime <= cleared_before:
                continue
            orphans.append({
                "id": "orphan_" + vd.name,
                "kind": "image" if "image" in reason else "text",
                "text": "(applied earlier — metadata unavailable)" if "text" in reason else None,
                "filename": "(image watermark — earlier session)" if "image" in reason else None,
                "scope": "remove" if "remove" in reason else "apply",
                "count": len(list(vd.glob("slide-*.jpg"))),
                "snapshot": vd.name,
                "timestamp": ts,
                "reason": reason,
                "orphan": True,
            })

    # Combine — newest first by timestamp
    # Hide entries that have been undone — they're no longer "currently applied".
    visible_entries = [e for e in entries if not e.get("undone")]
    combined = list(visible_entries) + orphans
    combined.sort(key=lambda e: e.get("timestamp", ""), reverse=True)

    out = []
    for e in combined:
        snap = e.get("snapshot")
        out.append({
            **e,
            "revertable": bool(snap and (HISTORY_DIR / snap).exists()),
        })
    return jsonify({"entries": out})


@app.route("/api/watermarks/revert/<entry_id>", methods=["POST"])
def revert_watermark(entry_id):
    """Revert a single watermark apply: restore its snapshot and drop the entry
    (plus any later entries, since reverting earlier invalidates them).
    Accepts both logged entry ids and orphan_<snapshot> ids for snapshots that
    exist on disk but were never logged (older session, cleared log, etc.)."""
    if not re.match(r"^[A-Za-z0-9_]{1,64}$", entry_id):
        return jsonify({"error": "Invalid id"}), 400

    entries = _load_wm_log()
    target_idx = None
    target = None
    snapshot = None

    if entry_id.startswith("orphan_"):
        snapshot = entry_id[len("orphan_"):]
        if not re.match(r"^[A-Za-z0-9_\-]{1,64}$", snapshot):
            return jsonify({"error": "Invalid orphan snapshot id"}), 400
        target = {"id": entry_id, "snapshot": snapshot, "orphan": True}
    else:
        for i, e in enumerate(entries):
            if e.get("id") == entry_id:
                target_idx = i
                target = e
                snapshot = e.get("snapshot")
                break
        if target_idx is None:
            return jsonify({"error": "Watermark entry not found"}), 404

    if not snapshot:
        return jsonify({"error": "No snapshot stored for this entry"}), 400

    version_dir = (HISTORY_DIR / snapshot).resolve()
    if HISTORY_DIR.resolve() not in version_dir.parents:
        return jsonify({"error": "Invalid snapshot path"}), 400
    if not version_dir.exists():
        return jsonify({"error": "Snapshot no longer exists"}), 404

    # ── Restore: scope-aware ────────────────────────────────────────────
    # If the watermark was applied to a single slide, restoring all slides
    # would wipe unrelated edits the user made afterwards. Honour the logged
    # scope/slide_num and restore only the affected slide in that case.
    is_current_scope = (
        target_idx is not None
        and target.get("scope") == "current"
        and target.get("slide_num")
    )
    data_file = version_dir / "data.json"

    if is_current_scope:
        n = int(target["slide_num"])
        name = f"slide-{n:02d}.jpg"
        src = version_dir / name
        dst = SLIDES_DIR / name
        if not src.exists():
            return jsonify({"error": "Slide missing from snapshot"}), 404
        shutil.copy2(str(src), str(dst))
        # Restore only this slide's overlays/notes from the snapshot's data.json.
        if data_file.exists():
            with _data_lock:
                live = json.loads(DATA_FILE.read_text()) if DATA_FILE.exists() else {}
                snap_data = json.loads(data_file.read_text())
                live[str(n)] = snap_data.get(str(n), {"overlays": [], "notes": ""})
                DATA_FILE.write_text(json.dumps(live, indent=2))
    else:
        # scope == "all" or orphan — restore everything (existing behaviour).
        for f in SLIDES_DIR.glob("slide-*.jpg"):
            f.unlink()
        for sf in sorted(version_dir.glob("slide-*.jpg")):
            shutil.copy2(str(sf), str(SLIDES_DIR / sf.name))
        if data_file.exists():
            with _data_lock:
                DATA_FILE.write_text(data_file.read_text())

    # ── Drop log entries that no longer match disk state ────────────────
    if target_idx is not None:
        if is_current_scope:
            # Only later entries that touched the SAME slide (or "all") are invalidated.
            target_slide = target.get("slide_num")
            kept = list(entries[:target_idx])
            dropped = [target]  # the reverted entry itself is always dropped
            for e in entries[target_idx + 1:]:
                touches_same = (
                    e.get("scope") == "all"
                    or (e.get("scope") == "current" and e.get("slide_num") == target_slide)
                )
                if touches_same:
                    dropped.append(e)
                else:
                    kept.append(e)
            _save_wm_log(kept)
        else:
            # scope=all revert invalidates this entry and everything after it.
            dropped = entries[target_idx:]
            _save_wm_log(entries[:target_idx])
    else:
        # Orphan revert: clear any log entries newer than the snapshot we just restored.
        snap_mtime = (HISTORY_DIR / snapshot).stat().st_mtime
        kept = [e for e in entries
                if not e.get("timestamp")
                or datetime.datetime.fromisoformat(e["timestamp"]).timestamp() < snap_mtime]
        dropped = entries[len(kept):]
        _save_wm_log(kept)
    return jsonify({"ok": True, "reverted": target, "dropped": len(dropped),
                    "slides": len(_get_slide_files())})


@app.route("/api/watermarks/clear-log", methods=["POST"])
def clear_watermark_log():
    """Clear the applied-watermark log AND hide any orphan snapshots created
    before now. Slide images and the underlying history snapshots are NOT
    touched — they remain available from the History modal."""
    import time as _t
    _save_wm_log([])
    _set_cleared_orphans_before(_t.time())
    return jsonify({"ok": True})


# ── Unified undo / redo across all server-side ops ──────────────────────────

@app.route("/api/ops/state", methods=["GET"])
def ops_state():
    """Return whether undo/redo are available, plus a one-line label for the
    next action in each direction. Used by the toolbar to enable/disable
    the Undo/Redo buttons and show a tooltip."""
    entries = _load_wm_log()
    last_active = None
    last_undone = None
    for e in entries:
        if e.get("undone"):
            last_undone = e
        else:
            last_active = e
    return jsonify({
        "can_undo": last_active is not None,
        "can_redo": last_undone is not None,
        "undo_label": (last_active and last_active.get("text")) or None,
        "redo_label": (last_undone and last_undone.get("text")) or None,
    })


@app.route("/api/ops/undo", methods=["POST"])
def ops_undo():
    """Undo the most recent server-side op that hasn't already been undone.

    Holds `_data_lock` (reentrant) across the entire sequence — find target,
    take redo snapshot, restore pre-op snapshot, mark entry — so concurrent
    requests can't interleave and corrupt the log.
    """
    with _data_lock:
        entries = _load_wm_log()
        target_idx = None
        for i in range(len(entries) - 1, -1, -1):
            if not entries[i].get("undone"):
                target_idx = i
                break
        if target_idx is None:
            return jsonify({"ok": False, "reason": "Nothing to undo"})

        target = entries[target_idx]
        snapshot_id = target.get("snapshot")
        if not snapshot_id:
            return jsonify({"error": "Entry has no snapshot"}), 400
        version_dir = (HISTORY_DIR / snapshot_id).resolve()
        if HISTORY_DIR.resolve() not in version_dir.parents:
            return jsonify({"error": "Invalid snapshot path"}), 400
        if not version_dir.exists():
            return jsonify({"error": "Snapshot no longer exists"}), 404

        try:
            redo_snapshot = _take_redo_snapshot()
            _restore_from_snapshot(version_dir, scope=target.get("scope"),
                                   slide_num=target.get("slide_num"))
        except (OSError, RuntimeError) as e:
            return jsonify({"error": f"Restore failed: {e}"}), 500

        entries[target_idx]["undone"] = True
        entries[target_idx]["redo_snapshot"] = redo_snapshot
        WATERMARK_LOG.write_text(json.dumps(entries, indent=2))

    return jsonify({"ok": True, "kind": target.get("kind"),
                    "text": target.get("text"), "id": target["id"]})


@app.route("/api/ops/redo", methods=["POST"])
def ops_redo():
    """Redo the most recently undone op by restoring its `redo_snapshot`.
    Holds `_data_lock` across the full sequence (see ops_undo)."""
    with _data_lock:
        entries = _load_wm_log()
        target_idx = None
        for i in range(len(entries) - 1, -1, -1):
            if entries[i].get("undone"):
                target_idx = i
                break
        if target_idx is None:
            return jsonify({"ok": False, "reason": "Nothing to redo"})
        target = entries[target_idx]
        redo_snap = target.get("redo_snapshot")
        if not redo_snap:
            return jsonify({"error": "Undone entry has no redo snapshot"}), 400
        version_dir = (HISTORY_DIR / redo_snap).resolve()
        if HISTORY_DIR.resolve() not in version_dir.parents:
            return jsonify({"error": "Invalid redo snapshot path"}), 400
        if not version_dir.exists():
            return jsonify({"error": "Redo snapshot no longer exists"}), 404

        try:
            _restore_from_snapshot(version_dir, scope=target.get("scope"),
                                   slide_num=target.get("slide_num"))
        except (OSError, RuntimeError) as e:
            return jsonify({"error": f"Restore failed: {e}"}), 500

        entries[target_idx]["undone"] = False
        entries[target_idx].pop("redo_snapshot", None)
        WATERMARK_LOG.write_text(json.dumps(entries, indent=2))

    return jsonify({"ok": True, "kind": target.get("kind"),
                    "text": target.get("text"), "id": target["id"]})


# ── Batch Processing ────────────────────────────────────────────────────────

@app.route("/api/batch/remove-logo", methods=["POST"])
def batch_remove_logo():
    """Upload up to 20 PPTX files, remove logos from all slides in each,
    and return a ZIP containing the cleaned PPTX files."""
    _cleanup_old_exports()
    import zipfile, uuid, tempfile

    files = request.files.getlist("files")
    if not files or not any(f.filename for f in files):
        return jsonify({"error": "No files uploaded"}), 400

    pptx_files = [f for f in files if f.filename and f.filename.lower().endswith(".pptx")]
    if not pptx_files:
        return jsonify({"error": "No .pptx files found"}), 400
    if len(pptx_files) > 20:
        return jsonify({"error": "Maximum 20 files allowed"}), 400

    tmp_dir = Path(tempfile.mkdtemp())
    results = []
    cleaned_paths = []

    try:
        for f in pptx_files:
            fname = secure_filename(f.filename)
            if not fname:
                continue
            input_path = tmp_dir / fname
            f.save(str(input_path))

            try:
                # Extract slides to a per-file temp dir
                file_slides_dir = tmp_dir / f"slides_{fname}"
                file_slides_dir.mkdir(exist_ok=True)

                # Convert PPTX to images (reuse shared conversion functions)
                _convert_pptx_to_images_libreoffice(input_path, file_slides_dir)

                # Remove logos from all slide images
                slide_images = sorted(file_slides_dir.glob("slide-*.jpg"))
                remove_logos_batch(slide_images)

                # Rebuild PPTX with cleaned images
                output_path = tmp_dir / f"clean_{fname}"
                _rebuild_pptx_from_images(slide_images, output_path)
                cleaned_paths.append((output_path, f"clean_{fname}"))
                results.append({"file": fname, "status": "ok", "slides": len(slide_images)})

            except Exception as e:
                results.append({"file": fname, "status": "error", "error": str(e)})

        if not cleaned_paths:
            return jsonify({"error": "No files were processed successfully", "results": results}), 500

        # Package all cleaned files into a ZIP
        zip_name = f"bulk_cleaned_{uuid.uuid4().hex[:8]}.zip"
        zip_path = EXPORT_DIR / zip_name
        with zipfile.ZipFile(str(zip_path), 'w', zipfile.ZIP_DEFLATED) as zf:
            for fpath, arcname in cleaned_paths:
                zf.write(str(fpath), arcname)

        return send_file(str(zip_path), as_attachment=True,
                         download_name="SlideCraft_Bulk_Cleaned.zip")

    finally:
        shutil.rmtree(str(tmp_dir), ignore_errors=True)


@app.route("/api/folder/remove-logo", methods=["POST"])
def folder_remove_logo():
    """Process every .pptx in a local folder: strip the NotebookLM logo and
    write the cleaned file to '<folder>/Slides Final/<name>.pptx'.

    Body: {"folder": "<abs path>", "overwrite": bool (optional, default False),
           "recursive": bool (optional, default False)}

    Streams NDJSON progress events as files complete:
      {"type":"start",  "total": N, "folder": ..., "output_folder": ...}
      {"type":"file",   "index": i, "total": N, "file": name,
                        "status": "ok"|"skipped"|"error", ...}
      {"type":"done",   "ok": K, "skipped": S, "error": E, "total": N}
    """
    import tempfile
    from flask import Response, stream_with_context

    payload = _ensure_dict(request.json)
    raw = payload.get("folder", "")
    if not isinstance(raw, str) or not raw.strip():
        return jsonify({"error": "Missing 'folder' path"}), 400

    src = Path(raw.strip()).expanduser()
    try:
        src = src.resolve(strict=True)
    except (OSError, RuntimeError):
        return jsonify({"error": f"Folder not found: {raw}"}), 400
    if not src.is_dir():
        return jsonify({"error": f"Not a folder: {src}"}), 400

    recursive = bool(payload.get("recursive", False))
    overwrite = bool(payload.get("overwrite", False))

    out_dir = src / "Slides Final"
    pattern = "**/*.pptx" if recursive else "*.pptx"
    pptx_files = sorted(
        p for p in src.glob(pattern)
        if p.is_file() and out_dir not in p.parents and p.parent != out_dir
    )
    if not pptx_files:
        return jsonify({"error": "No .pptx files found in folder",
                        "folder": str(src)}), 404

    out_dir.mkdir(exist_ok=True)
    total = len(pptx_files)

    def _emit(obj):
        return (json.dumps(obj) + "\n").encode("utf-8")

    @stream_with_context
    def generate():
        ok = skipped = errored = 0
        yield _emit({"type": "start", "total": total,
                     "folder": str(src), "output_folder": str(out_dir)})

        for i, input_path in enumerate(pptx_files, start=1):
            fname = input_path.name
            output_path = out_dir / fname

            yield _emit({"type": "file", "index": i, "total": total,
                         "file": fname, "status": "processing"})

            if output_path.exists() and not overwrite:
                skipped += 1
                yield _emit({"type": "file", "index": i, "total": total,
                             "file": fname, "status": "skipped",
                             "reason": "already exists in Slides Final"})
                continue

            tmp_dir = Path(tempfile.mkdtemp(prefix="slidecraft_folder_"))
            try:
                file_slides_dir = tmp_dir / "slides"
                file_slides_dir.mkdir()
                _convert_pptx_to_images_libreoffice(input_path, file_slides_dir)

                slide_images = sorted(file_slides_dir.glob("slide-*.jpg"))
                if not slide_images:
                    errored += 1
                    yield _emit({"type": "file", "index": i, "total": total,
                                 "file": fname, "status": "error",
                                 "error": "Could not render slides"})
                    continue

                remove_logos_batch(slide_images)
                _rebuild_pptx_from_images(slide_images, output_path)
                ok += 1
                yield _emit({"type": "file", "index": i, "total": total,
                             "file": fname, "status": "ok",
                             "slides": len(slide_images),
                             "output": str(output_path)})
            except Exception as e:
                errored += 1
                yield _emit({"type": "file", "index": i, "total": total,
                             "file": fname, "status": "error",
                             "error": str(e)})
            finally:
                shutil.rmtree(str(tmp_dir), ignore_errors=True)

        yield _emit({"type": "done", "total": total,
                     "ok": ok, "skipped": skipped, "error": errored,
                     "folder": str(src), "output_folder": str(out_dir)})

    return Response(generate(), mimetype="application/x-ndjson",
                    headers={"X-Accel-Buffering": "no",
                             "Cache-Control": "no-cache"})


def _rebuild_pptx_from_images(slide_images, output_path):
    """Rebuild a PPTX with slide images as full-slide backgrounds."""
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]
    for sf in slide_images:
        slide = prs.slides.add_slide(blank_layout)
        pic = slide.shapes.add_picture(
            str(sf), left=0, top=0,
            width=prs.slide_width, height=prs.slide_height)
        slide.shapes._spTree.remove(pic._element)
        slide.shapes._spTree.insert(2, pic._element)
    prs.save(str(output_path))


# ── Templates ───────────────────────────────────────────────────────────────

TEMPLATES_DIR = BASE_DIR / "templates_saved"
TEMPLATES_DIR.mkdir(exist_ok=True)


@app.route("/api/templates", methods=["GET"])
def list_templates():
    """List saved slide templates."""
    templates = []
    for tf in sorted(TEMPLATES_DIR.glob("*.json")):
        data = json.loads(tf.read_text())
        templates.append({"name": tf.stem, "slides": data.get("slide_count", 0),
                          "created": data.get("created", "")})
    return jsonify({"templates": templates})


@app.route("/api/templates/save", methods=["POST"])
def save_template():
    """Save current slides + overlays as a named template."""
    payload = _ensure_dict(request.json)
    name = _safe_name(payload.get("name", "untitled"))
    if not name:
        return jsonify({"error": "Invalid template name (use letters, numbers, _, -, space; max 64)"}), 400

    slide_files = _get_slide_files()
    data = load_data()

    # Copy slide images to template dir
    tpl_dir = TEMPLATES_DIR / name
    tpl_dir.mkdir(exist_ok=True)
    for sf in slide_files:
        shutil.copy2(str(sf), str(tpl_dir / sf.name))

    # Save overlay data
    meta = {"slide_count": len(slide_files), "data": data,
            "created": datetime.datetime.now().isoformat()}
    (TEMPLATES_DIR / f"{name}.json").write_text(json.dumps(meta, indent=2))

    return jsonify({"ok": True, "name": name})


@app.route("/api/templates/load", methods=["POST"])
def load_template():
    """Load a template — restore slide images + overlays."""
    payload = _ensure_dict(request.json)
    name = _safe_name(payload.get("name", ""))
    if not name:
        return jsonify({"error": "Invalid template name"}), 400

    tpl_dir = TEMPLATES_DIR / name
    meta_file = TEMPLATES_DIR / f"{name}.json"

    if not meta_file.exists():
        return jsonify({"error": "Template not found"}), 404

    snapshot = _snapshot_before_destructive("load-template")

    # Clear current slides
    for f in SLIDES_DIR.glob("slide-*.jpg"):
        f.unlink()

    # Copy template slides
    for sf in sorted(tpl_dir.glob("slide-*.jpg")):
        shutil.copy2(str(sf), str(SLIDES_DIR / sf.name))

    # Restore overlay data
    meta = json.loads(meta_file.read_text())
    save_data(meta.get("data", {}))

    slide_count = len(_get_slide_files())
    log_id = _log_op("load-template", text=f"Loaded template '{name}'",
                     scope="all", count=slide_count, snapshot=snapshot)
    return jsonify({"ok": True, "slides": slide_count,
                    "snapshot": snapshot, "log_id": log_id})


@app.route("/api/templates/delete", methods=["POST"])
def delete_template():
    """Delete a saved template."""
    payload = _ensure_dict(request.json)
    name = _safe_name(payload.get("name", ""))
    if not name:
        return jsonify({"error": "Invalid template name"}), 400
    tpl_dir = (TEMPLATES_DIR / name).resolve()
    # Confine to TEMPLATES_DIR
    if TEMPLATES_DIR.resolve() not in tpl_dir.parents and tpl_dir != TEMPLATES_DIR.resolve():
        return jsonify({"error": "Invalid path"}), 400
    meta_file = TEMPLATES_DIR / f"{name}.json"
    if tpl_dir.exists() and tpl_dir != TEMPLATES_DIR.resolve():
        shutil.rmtree(str(tpl_dir))
    meta_file.unlink(missing_ok=True)
    return jsonify({"ok": True})


# ── Version History ─────────────────────────────────────────────────────────

HISTORY_DIR = BASE_DIR / "history"
HISTORY_DIR.mkdir(exist_ok=True)


@app.route("/api/history/save", methods=["POST"])
def save_version():
    """Save a snapshot of current slide state."""
    ts = datetime.datetime.now().strftime("%Y%m%d_%H%M%S") + f"_{uuid.uuid4().hex[:4]}"
    version_dir = HISTORY_DIR / ts

    version_dir.mkdir(exist_ok=True)
    for sf in _get_slide_files():
        shutil.copy2(str(sf), str(version_dir / sf.name))

    data = load_data()
    (version_dir / "data.json").write_text(json.dumps(data, indent=2))

    return jsonify({"ok": True, "version": ts})


@app.route("/api/history", methods=["GET"])
def list_versions():
    """List all saved versions."""
    versions = []
    for vd in sorted(HISTORY_DIR.iterdir(), reverse=True):
        if vd.is_dir():
            slides = len(list(vd.glob("slide-*.jpg")))
            versions.append({"version": vd.name, "slides": slides})
    return jsonify({"versions": versions})


@app.route("/api/history/restore", methods=["POST"])
def restore_version():
    """Restore a previous version."""
    payload = _ensure_dict(request.json)
    raw = payload.get("version", "") or payload.get("name", "")
    # Version names are timestamps + hex, so allow that pattern strictly.
    if not isinstance(raw, str) or not re.match(r"^[A-Za-z0-9_\-]{1,64}$", raw):
        return jsonify({"error": "Invalid version"}), 400
    version_dir = (HISTORY_DIR / raw).resolve()
    if HISTORY_DIR.resolve() not in version_dir.parents:
        return jsonify({"error": "Invalid path"}), 400
    if not version_dir.exists() or version_dir == HISTORY_DIR.resolve():
        return jsonify({"error": "Version not found"}), 404

    snapshot = _snapshot_before_destructive("restore-version")
    for f in SLIDES_DIR.glob("slide-*.jpg"):
        f.unlink()
    for sf in sorted(version_dir.glob("slide-*.jpg")):
        shutil.copy2(str(sf), str(SLIDES_DIR / sf.name))

    data_file = version_dir / "data.json"
    if data_file.exists():
        save_data(json.loads(data_file.read_text()))

    slide_count = len(_get_slide_files())
    log_id = _log_op("restore-version", text=f"Restored version '{raw}'",
                     scope="all", count=slide_count, snapshot=snapshot)
    return jsonify({"ok": True, "slides": slide_count,
                    "snapshot": snapshot, "log_id": log_id})


# ── Comments / Annotations ──────────────────────────────────────────────────

COMMENTS_FILE = BASE_DIR / "comments.json"


def _load_comments():
    with _data_lock:
        if COMMENTS_FILE.exists():
            return json.loads(COMMENTS_FILE.read_text())
        return {}


def _save_comments(data):
    with _data_lock:
        COMMENTS_FILE.write_text(json.dumps(data, indent=2))


@app.route("/api/comments/<int:num>", methods=["GET"])
def get_comments(num):
    data = _load_comments()
    return jsonify({"comments": data.get(str(num), [])})


@app.route("/api/comments/<int:num>", methods=["POST"])
def add_comment(num):
    payload = _ensure_dict(request.json)
    comment = {
        "text": str(payload.get("text", ""))[:5000],
        "x": float(payload.get("x", 0.5)) if isinstance(payload.get("x", 0.5), (int, float)) else 0.5,
        "y": float(payload.get("y", 0.5)) if isinstance(payload.get("y", 0.5), (int, float)) else 0.5,
        "author": str(payload.get("author", "User"))[:80],
        "timestamp": datetime.datetime.now().isoformat(),
        "resolved": False,
    }
    with _data_lock:
        data = json.loads(COMMENTS_FILE.read_text()) if COMMENTS_FILE.exists() else {}
        data.setdefault(str(num), []).append(comment)
        COMMENTS_FILE.write_text(json.dumps(data, indent=2))
    return jsonify({"ok": True, "comment": comment})


@app.route("/api/comments/<int:num>/resolve", methods=["POST"])
def resolve_comment(num):
    payload = _ensure_dict(request.json)
    try:
        idx = int(payload.get("index", 0))
    except (TypeError, ValueError):
        return jsonify({"error": "Invalid index"}), 400
    with _data_lock:
        data = json.loads(COMMENTS_FILE.read_text()) if COMMENTS_FILE.exists() else {}
        comments = data.get(str(num), [])
        if 0 <= idx < len(comments):
            comments[idx]["resolved"] = True
            COMMENTS_FILE.write_text(json.dumps(data, indent=2))
    return jsonify({"ok": True})


@app.route("/api/comments/<int:num>/delete", methods=["POST"])
def delete_comment(num):
    payload = _ensure_dict(request.json)
    try:
        idx = int(payload.get("index", 0))
    except (TypeError, ValueError):
        return jsonify({"error": "Invalid index"}), 400
    with _data_lock:
        data = json.loads(COMMENTS_FILE.read_text()) if COMMENTS_FILE.exists() else {}
        comments = data.get(str(num), [])
        if 0 <= idx < len(comments):
            comments.pop(idx)
            COMMENTS_FILE.write_text(json.dumps(data, indent=2))
    return jsonify({"ok": True})


# ── Find & Replace ──────────────────────────────────────────────────────────

@app.route("/api/find-replace", methods=["POST"])
def find_replace():
    """Find and replace text in overlays across all slides."""
    payload = _ensure_dict(request.json)
    find_text = str(payload.get("find", ""))
    replace_text = str(payload.get("replace", ""))
    if not find_text:
        return jsonify({"error": "No search text"}), 400

    snapshot = _snapshot_before_destructive("find-replace")
    count = 0
    with _data_lock:
        data = json.loads(DATA_FILE.read_text()) if DATA_FILE.exists() else {}
        for slide_num, slide_data in data.items():
            for ov in slide_data.get("overlays", []):
                if ov.get("type") == "text" and find_text in ov.get("text", ""):
                    ov["text"] = ov["text"].replace(find_text, replace_text)
                    count += 1
        DATA_FILE.write_text(json.dumps(data, indent=2))
    log_id = _log_op("find-replace",
                     text=f"Replaced {count} '{find_text}' → '{replace_text}'",
                     scope="all", count=count, snapshot=snapshot)
    return jsonify({"ok": True, "replacements": count,
                    "snapshot": snapshot, "log_id": log_id})


# ── Video Logo Removal ──────────────────────────────────────────────────────

@app.route("/video")
def video_page():
    return render_template("video.html")


@app.route("/api/video/upload", methods=["POST"])
def upload_video():
    """Upload a video file for logo removal."""
    if "file" not in request.files:
        return jsonify({"error": "No file"}), 400
    f = request.files["file"]
    fname = secure_filename(f.filename)
    if not fname:
        return jsonify({"error": "Invalid filename"}), 400
    ext = Path(fname).suffix.lower()
    if ext not in ALLOWED_VIDEO_EXTS:
        return jsonify({"error": f"Unsupported video format: {ext}. Allowed: {', '.join(ALLOWED_VIDEO_EXTS)}"}), 400

    save_path = VIDEO_DIR / fname
    f.save(str(save_path))
    return jsonify({"ok": True, "filename": fname})


@app.route("/api/video/preview-frame", methods=["POST"])
def video_preview_frame():
    """Get a frame from the video + detect logo region."""
    payload = request.json
    fname = payload.get("filename", "")
    video_path = VIDEO_DIR / secure_filename(fname)
    if not video_path.exists():
        return jsonify({"error": "Video not found"}), 400

    import cv2
    cap = cv2.VideoCapture(str(video_path))
    cap.set(cv2.CAP_PROP_POS_FRAMES, 30)  # grab frame 30
    ret, frame = cap.read()
    w = int(cap.get(cv2.CAP_PROP_FRAME_WIDTH))
    h = int(cap.get(cv2.CAP_PROP_FRAME_HEIGHT))
    fps = cap.get(cv2.CAP_PROP_FPS)
    total = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
    duration = total / fps if fps > 0 else 0
    cap.release()

    if not ret:
        return jsonify({"error": "Could not read video frame"}), 500

    # Convert frame to base64 JPEG
    _, buf = cv2.imencode('.jpg', frame, [cv2.IMWRITE_JPEG_QUALITY, 85])
    b64 = base64.b64encode(buf.tobytes()).decode()

    return jsonify({
        "frame": f"data:image/jpeg;base64,{b64}",
        "width": w, "height": h,
        "fps": fps, "duration": round(duration, 1),
        "totalFrames": total
    })


import time as _time

_video_jobs = {}  # job_id → {status, progress, total, eta, output, error, created_at}
_VIDEO_JOB_MAX_AGE = 3600  # expire jobs after 1 hour


def _cleanup_video_jobs():
    """Remove completed/cancelled jobs older than max age."""
    now = _time.time()
    expired = [jid for jid, job in _video_jobs.items()
               if job.get("status") in ("done", "cancelled", "error")
               and now - job.get("created_at", 0) > _VIDEO_JOB_MAX_AGE]
    for jid in expired:
        del _video_jobs[jid]


@app.route("/api/video/remove-logo", methods=["POST"])
def remove_video_logo():
    """Start logo removal job in a background thread, return job_id for polling."""
    payload = request.json
    fname = payload.get("filename", "")
    lx = float(payload.get("x", 0.85))
    ly = float(payload.get("y", 0.93))
    lw = float(payload.get("w", 0.14))
    lh = float(payload.get("h", 0.06))

    video_path = VIDEO_DIR / secure_filename(fname)
    if not video_path.exists():
        return jsonify({"error": "Video not found"}), 400

    _cleanup_video_jobs()
    job_id = f"job_{int(_time.time()*1000)}"
    _video_jobs[job_id] = {"status": "starting", "progress": 0, "total": 0, "eta": 0, "output": "", "error": "", "created_at": _time.time()}

    t = threading.Thread(target=_run_logo_removal, args=(job_id, video_path, fname, lx, ly, lw, lh), daemon=True)
    t.start()

    return jsonify({"ok": True, "jobId": job_id})


@app.route("/api/video/progress/<job_id>")
def video_progress(job_id):
    """Poll for job progress."""
    job = _video_jobs.get(job_id)
    if not job:
        return jsonify({"error": "Job not found"}), 404
    return jsonify(job)


@app.route("/api/video/stop/<job_id>", methods=["POST"])
def stop_video_job(job_id):
    """Request cancellation of a running video job."""
    job = _video_jobs.get(job_id)
    if not job:
        return jsonify({"error": "Job not found"}), 404
    job["cancel"] = True
    return jsonify({"ok": True})


def _run_logo_removal(job_id, video_path, fname, lx, ly, lw, lh):
    """Background worker for video logo removal."""
    import cv2
    import numpy as np

    job = _video_jobs[job_id]
    job["status"] = "processing"

    try:
        cap = cv2.VideoCapture(str(video_path))
        w = int(cap.get(cv2.CAP_PROP_FRAME_WIDTH))
        h = int(cap.get(cv2.CAP_PROP_FRAME_HEIGHT))
        fps = cap.get(cv2.CAP_PROP_FPS)
        total = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
        job["total"] = total

        x1 = int(lx * w)
        y1 = int(ly * h)
        x2 = min(w, int((lx + lw) * w))
        y2 = min(h, int((ly + lh) * h))

        out_name = f"clean_{secure_filename(fname)}"
        if not out_name.lower().endswith('.mp4'):
            out_name = out_name.rsplit('.', 1)[0] + '.mp4'
        out_path = VIDEO_DIR / out_name
        temp_path = VIDEO_DIR / f"_temp_{out_name}"

        fourcc = cv2.VideoWriter_fourcc(*'mp4v')
        writer = cv2.VideoWriter(str(temp_path), fourcc, fps, (w, h))

        mask = np.zeros((h, w), dtype=np.uint8)
        mask[y1:y2, x1:x2] = 255

        frame_num = 0
        start_time = _time.time()

        while True:
            ret, frame = cap.read()
            if not ret:
                break

            # Check for cancellation
            if job.get("cancel"):
                job["status"] = "cancelled"
                break

            frame = cv2.inpaint(frame, mask, inpaintRadius=5, flags=cv2.INPAINT_TELEA)
            writer.write(frame)
            frame_num += 1

            # Update progress every 10 frames
            if frame_num % 10 == 0 or frame_num == total:
                elapsed = _time.time() - start_time
                fps_actual = frame_num / max(0.1, elapsed)
                remaining = (total - frame_num) / max(1, fps_actual)
                job["progress"] = frame_num
                job["eta"] = round(remaining, 1)

        cap.release()
        writer.release()

        # If cancelled, clean up and exit
        if job.get("cancel"):
            temp_path.unlink(missing_ok=True)
            job["eta"] = 0
            return

        # Mux audio
        job["status"] = "muxing_audio"
        job["eta"] = 0
        try:
            from moviepy import VideoFileClip
            original = VideoFileClip(str(video_path))
            clean = VideoFileClip(str(temp_path))
            if original.audio is not None:
                final = clean.with_audio(original.audio)
                temp_audio = VIDEO_DIR / f'_temp_audio_{job_id}.m4a'
                final.write_videofile(str(out_path), codec='libx264', audio_codec='aac',
                                      logger=None, temp_audiofile=str(temp_audio))
                final.close()
            else:
                clean.write_videofile(str(out_path), codec='libx264', logger=None)
            original.close()
            clean.close()
            temp_path.unlink(missing_ok=True)
        except ImportError:
            if temp_path.exists():
                shutil.move(str(temp_path), str(out_path))
            job["warning"] = "moviepy not installed — audio track not restored"
        except Exception as _mux_err:
            if temp_path.exists():
                shutil.move(str(temp_path), str(out_path))
            job["warning"] = f"Audio muxing failed, output is video-only: {_mux_err}"
            print(f"[video-mux] {_mux_err}", file=sys.stderr)

        (VIDEO_DIR / f'_temp_audio_{job_id}.m4a').unlink(missing_ok=True)

        job["status"] = "done"
        job["progress"] = total
        job["output"] = out_name
        job["eta"] = 0

    except Exception as e:
        job["status"] = "error"
        job["error"] = str(e)


@app.route("/api/video/download/<filename>")
def download_video(filename):
    safe_name = secure_filename(filename)
    fpath = VIDEO_DIR / safe_name
    if not fpath.exists():
        return jsonify({"error": "File not found"}), 404
    return send_file(str(fpath), as_attachment=True, download_name=safe_name)


# ── Feature modules (registered after all core routes) ──────────────────────
_FEATURE_CTX = {
    "BASE_DIR":     BASE_DIR,
    "SLIDES_DIR":   SLIDES_DIR,
    "ORIGINALS_DIR": ORIGINALS_DIR,
    "DATA_FILE":    DATA_FILE,
    "UPLOAD_DIR":   UPLOAD_DIR,
    "EXPORT_DIR":   EXPORT_DIR,
    "VIDEO_DIR":    VIDEO_DIR,
    "load_data":    load_data,
    "save_data":    save_data,
    "_get_slide_files": _get_slide_files,
    "_safe_name":   _safe_name,
    "_data_lock":   _data_lock,
    "_set_deck_name": _set_deck_name,
    "_get_deck_name": _get_deck_name,
    "MAX_OVERLAY_IMG_BYTES": MAX_OVERLAY_IMG_BYTES,
}

from app_features import register_feature_routes
register_feature_routes(app, _FEATURE_CTX)

from app_ai import register_ai_routes
register_ai_routes(app, _FEATURE_CTX)

from app_auth import register_auth_routes
register_auth_routes(app, _FEATURE_CTX)

# Google Slides integration is optional — requires OAuth credentials in env.
try:
    from app_gslides import register_gslides_routes
    register_gslides_routes(app, _FEATURE_CTX)
except ImportError as _e:
    print(f"[gslides] optional module not loaded: {_e}", file=sys.stderr)


def _clear_session():
    """Wipe slide images and overlay data so every startup is a clean slate."""
    for f in SLIDES_DIR.glob("slide-*.jpg"):
        try: f.unlink()
        except OSError: pass
    for path in (DATA_FILE, DECK_NAME_FILE, BASE_DIR / "autosave.json"):
        try: path.unlink()
        except OSError: pass


if __name__ == "__main__":
    if not _find_libreoffice():
        print(
            "ERROR: LibreOffice not found. SlideCraft requires LibreOffice for PPTX conversion.\n"
            "Install from https://www.libreoffice.org/download/download/ and restart.",
            file=sys.stderr,
        )
        sys.exit(1)
    _clear_session()
    # Bind to localhost by default. Set HOST=0.0.0.0 to expose on LAN (no auth!).
    host = os.environ.get('HOST', '127.0.0.1')
    if host == '0.0.0.0':
        print("WARNING: binding to 0.0.0.0 — anyone on your network can read, edit, or delete your slides.", file=sys.stderr)
    app.run(debug=os.environ.get('FLASK_DEBUG', 'false').lower() == 'true',
            port=int(os.environ.get('PORT', 5050)), host=host)
