"""Shared fixtures. Each test gets an isolated working directory so that
state files (slide_data.json, comments.json, history/, etc.) don't leak."""
import sys
from pathlib import Path
import pytest
from PIL import Image

ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ROOT))

# Import app exactly ONCE for the whole test session. Re-importing per test
# pulls cv2 + easyocr + torch + moviepy fresh each time and adds ~150ms per
# test. We isolate state by rebinding module-level path attributes inside the
# per-test fixture instead.
import app as _app_module  # noqa: E402

_ORIGINAL_PATHS = {
    "BASE_DIR": _app_module.BASE_DIR,
    "SLIDES_DIR": _app_module.SLIDES_DIR,
    "ORIGINALS_DIR": _app_module.ORIGINALS_DIR,
    "UPLOAD_DIR": _app_module.UPLOAD_DIR,
    "EXPORT_DIR": _app_module.EXPORT_DIR,
    "DATA_FILE": _app_module.DATA_FILE,
    "PDF_TEXT_FILE": _app_module.PDF_TEXT_FILE,
    "VIDEO_DIR": _app_module.VIDEO_DIR,
    "TEMPLATES_DIR": _app_module.TEMPLATES_DIR,
    "HISTORY_DIR": _app_module.HISTORY_DIR,
    "COMMENTS_FILE": _app_module.COMMENTS_FILE,
    "WATERMARK_LOG": _app_module.WATERMARK_LOG,
    "CLEARED_ORPHANS_MARKER": _app_module.CLEARED_ORPHANS_MARKER,
}


def _make_dummy_pptx(path: Path, num_slides: int = 2):
    """Build a minimal valid PPTX with `num_slides` blank slides."""
    from pptx import Presentation
    prs = Presentation()
    blank = prs.slide_layouts[6]
    for _ in range(num_slides):
        prs.slides.add_slide(blank)
    prs.save(str(path))


def _make_dummy_jpg(path: Path, w: int = 800, h: int = 450, color=(220, 220, 220)):
    Image.new("RGB", (w, h), color).save(str(path), "JPEG", quality=90)


@pytest.fixture
def app_client(tmp_path, monkeypatch):
    """Flask test client with isolated state directories.

    Reuses the single session-wide `app` import; only the module-level paths
    are rebound to point inside `tmp_path` so each test gets a clean slate.
    """
    monkeypatch.chdir(tmp_path)

    app_module = _app_module
    # Rebind every writable path to live under tmp_path for this test only.
    app_module.BASE_DIR = tmp_path
    app_module.SLIDES_DIR = tmp_path / "static" / "slides"
    app_module.ORIGINALS_DIR = app_module.SLIDES_DIR / "_originals"
    app_module.UPLOAD_DIR = tmp_path / "uploads"
    app_module.EXPORT_DIR = tmp_path / "exports"
    app_module.DATA_FILE = tmp_path / "slide_data.json"
    app_module.PDF_TEXT_FILE = tmp_path / "pdf_text.json"
    app_module.VIDEO_DIR = tmp_path / "videos"
    app_module.TEMPLATES_DIR = tmp_path / "templates_saved"
    app_module.HISTORY_DIR = tmp_path / "history"
    app_module.COMMENTS_FILE = tmp_path / "comments.json"
    app_module.WATERMARK_LOG = tmp_path / "watermarks_applied.json"
    app_module.CLEARED_ORPHANS_MARKER = tmp_path / "cleared_orphans_before.txt"

    for d in (app_module.SLIDES_DIR, app_module.ORIGINALS_DIR, app_module.UPLOAD_DIR,
              app_module.EXPORT_DIR, app_module.VIDEO_DIR, app_module.TEMPLATES_DIR,
              app_module.HISTORY_DIR):
        d.mkdir(parents=True, exist_ok=True)

    app_module.app.config["TESTING"] = True
    client = app_module.app.test_client()

    yield client, app_module

    # Restore the original module-level paths so other tests / the live server
    # don't see tmp paths leak through.
    for name, value in _ORIGINAL_PATHS.items():
        setattr(app_module, name, value)


@pytest.fixture
def app_with_slides(app_client):
    """Pre-populated with 3 dummy slide JPGs."""
    client, app_module = app_client
    for i in range(1, 4):
        name = f"slide-{i:03d}.jpg"
        _make_dummy_jpg(app_module.SLIDES_DIR / name)
        _make_dummy_jpg(app_module.ORIGINALS_DIR / name)
    return client, app_module


@pytest.fixture
def dummy_pptx(tmp_path):
    p = tmp_path / "fixture.pptx"
    _make_dummy_pptx(p, 2)
    return p
